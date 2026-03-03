from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def bullet_list(s, l, t, w, h, items, sz=13, color=LIGHT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = False
        p.space_after = Pt(6)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === HEADER ===
circ(slide, 0.5, 0.2, 0.55, ORANGE, None)
txt(slide, 0.5, 0.24, 0.55, 0.45, "4", 22, DARK, True, PP_ALIGN.CENTER)
txt(slide, 1.2, 0.2, 12, 0.6, "STREAM PROCESSING", 36, ORANGE, True)
txt(slide, 1.2, 0.75, 14, 0.4,
    "Amazon Managed Service for Apache Flink — Real-time analytics on vehicle data", 22, ORANGE, True)

# === LEFT COLUMN — Capabilities ===
lx = 0.6
txt(slide, lx, 1.4, 7, 0.35, "Capabilities", 18, WHITE, True)

rect(slide, lx, 1.85, 7.0, 4.6, CARD, ORANGE, 1.5)

bullet_list(slide, lx + 0.25, 1.95, 6.5, 4.4, [
    ("Fully managed Apache Flink for real-time processing of", WHITE),
    ("vehicle telemetry streams with sub-second latency", LIGHT),
    ("", LIGHT),
    ("▸  Stateful stream processing — windowed aggregations,", LIGHT),
    ("   pattern detection, and sessionization across vehicle events", LIGHT),
    ("", LIGHT),
    ("▸  Real-time anomaly detection — identify battery degradation,", LIGHT),
    ("   sensor drift, or driving behavior changes as they happen", LIGHT),
    ("", LIGHT),
    ("▸  Complex event processing (CEP) — correlate signals across", LIGHT),
    ("   multiple ECUs to detect multi-signal failure patterns", LIGHT),
    ("", LIGHT),
    ("▸  Native MSK/Kafka integration — consume directly from", LIGHT),
    ("   Kafka topics, write results to S3, Timestream, or Kinesis", LIGHT),
    ("", LIGHT),
    ("▸  SQL and Java/Python APIs — data engineers use Flink SQL,", LIGHT),
    ("   platform teams use Java/Python for complex logic", LIGHT),
], 13, LIGHT)

# === RIGHT COLUMN — Differentiators ===
rx = 7.9
txt(slide, rx, 1.4, 7, 0.35, "Differentiators", 18, WHITE, True)

rect(slide, rx, 1.85, 7.5, 4.6, CARD, YELLOW, 1.5)

bullet_list(slide, rx + 0.25, 1.95, 7.0, 4.4, [
    ("vs. Databricks Structured Streaming", WHITE),
    ("Flink delivers true sub-second latency for safety-critical", LIGHT),
    ("use cases (DTC alerts, geofencing). Spark Structured Streaming", LIGHT),
    ("is micro-batch — seconds to minutes of latency", LIGHT),
    ("", LIGHT),
    ("vs. Decodable / Striim", WHITE),
    ("Managed Flink offers full Flink API access for complex stateful", LIGHT),
    ("logic. Decodable/Striim are simpler but limited for advanced", LIGHT),
    ("CEP patterns needed in vehicle anomaly detection", LIGHT),
    ("", LIGHT),
    ("vs. Self-managed Flink", WHITE),
    ("Zero operational overhead — auto-scaling, checkpointing,", LIGHT),
    ("snapshots, and upgrades are fully managed. No clusters to", LIGHT),
    ("provision, patch, or right-size", LIGHT),
    ("", LIGHT),
    ("Pay-as-you-go", WHITE),
    ("No upfront costs or licenses — pay only for KPUs consumed", LIGHT),
], 13, LIGHT)

# === BOTTOM — Automotive use cases ===
py = 6.7
txt(slide, 0.6, py, 4, 0.3, "Automotive Use Cases", 14, MUTED, True)

products = [
    ("Anomaly Detection", ORANGE, "Battery degradation, sensor drift, DTC correlation"),
    ("Geofencing & Alerts", GREEN, "Real-time location triggers and driver notifications"),
    ("Predictive Signals", PURPLE, "Enrich telemetry with ML features for downstream models"),
]

px = 0.6
for name, color, desc in products:
    rect(slide, px, py + 0.35, 4.8, 0.7, CARD, color, 1)
    txt(slide, px + 0.2, py + 0.38, 2.5, 0.3, name, 14, WHITE, True)
    txt(slide, px + 0.2, py + 0.65, 4.4, 0.3, desc, 11, LIGHT)
    px += 5.0

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide13_flink.pptx")
prs.save(out)
print(f"Saved: {out}")
