from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def pill(s, l, t, text, color, sz=11):
    pw = len(text) * 0.085 + 0.3
    ph = 0.28
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(pw), Inches(ph))
    dark = RGBColor(color[0] // 6, color[1] // 6, color[2] // 6)
    sh.fill.solid(); sh.fill.fore_color.rgb = dark
    sh.line.color.rgb = color; sh.line.width = Pt(1)
    tb = sh.text_frame; tb.word_wrap = False
    p = tb.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz); p.font.color.rgb = color
    return pw

# === TITLE ===
txt(slide, 0.6, 0.2, 14, 0.6, "PARTNERS THAT COMPLETE THE PLATFORM", 36, ORANGE, True)
txt(slide, 0.6, 0.75, 14, 0.4,
    "Specialized capabilities that extend the connected vehicle ecosystem", 20, LIGHT)

# === CAPABILITY CARDS — 4 columns x 2 rows ===
capabilities = [
    # Row 1
    ("Navigation &\nHD Maps", SKY,
     "HERE Technologies",
     "$1B AWS deal · SDV Accelerator",
     "AWS Service: Amazon Location Service"),
    ("OTA Updates &\nDevice Mgmt", GREEN,
     "Excelfore",
     "20M+ vehicles · AWS Competency",
     "AWS Service: AWS IoT Jobs"),
    ("Vehicle\nCybersecurity", RED,
     "Upstream Security",
     "XDR platform · UNECE R155",
     "AWS Service: AWS IoT Device Defender"),
    ("Predictive\nMaintenance", YELLOW,
     "Intangles",
     "96% fault prediction · Digital twin",
     "AWS Service: Amazon Lookout for Equipment"),
    # Row 2
    ("PKI & Certificate\nManagement", PURPLE,
     "Digicert",
     "Automotive-grade V2X PKI",
     "AWS Service: AWS Private CA"),
    ("Digital Twin &\nSimulation", PINK,
     "dSPACE",
     "AWS Competency · HIL + cloud sim",
     "AWS Service: AWS IoT TwinMaker"),
    ("Content Delivery\nfor OTA", SKY,
     "AWS CloudFront",
     "Global edge · signed URLs",
     "Native AWS — no partner needed"),
    ("Fleet Management\n& Telematics", GREEN,
     "Geotab",
     "4M+ vehicles · open platform",
     "AWS Service: AWS IoT FleetWise"),
]

cw = 3.55
ch = 2.2
cgap_x = 0.2
cgap_y = 0.2
start_x = 0.6
start_y = 1.4

for i, (name, color, partner, detail, aws_svc) in enumerate(capabilities):
    col = i % 4
    row = i // 4
    x = start_x + col * (cw + cgap_x)
    y = start_y + row * (ch + cgap_y)

    # Card
    rect(slide, x, y, cw, ch, CARD, color, 1.5)

    # Capability name
    txt(slide, x + 0.2, y + 0.12, cw - 0.4, 0.55, name, 14, WHITE, True)

    # Partner name
    txt(slide, x + 0.2, y + 0.72, cw - 0.4, 0.3, partner, 13, color, True)

    # Detail
    txt(slide, x + 0.2, y + 1.02, cw - 0.4, 0.35, detail, 10, LIGHT)

    # AWS service line
    txt(slide, x + 0.2, y + 1.42, cw - 0.4, 0.55, aws_svc, 9, MUTED)

# === BOTTOM BAR — key message ===
bar_y = start_y + 2 * (ch + cgap_y) + 0.15
rect(slide, 0.6, bar_y, 14.8, 0.5, CARD, ORANGE, 1.5)
txt(slide, 0.9, bar_y + 0.07, 6, 0.35,
    "AWS provides the foundation — partners fill the gaps", 15, ORANGE, True)
txt(slide, 7.5, bar_y + 0.1, 7.5, 0.3,
    "Every capability runs on AWS infrastructure · Available on AWS Marketplace", 12, LIGHT)

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide9_partners_gaps.pptx")
prs.save(out)
print(f"Saved: {out}")
