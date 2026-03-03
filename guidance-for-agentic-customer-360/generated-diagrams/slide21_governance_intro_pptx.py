from pptx import Presentation
import os

TEMPLATE = '/Users/givenand/Downloads/CVx-FCD-2026.pptx.pptx'
OUTDIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation(TEMPLATE)
slide = prs.slides[0]

updates = {
    'TextBox 10': 'Data Mesh Foundation: The solution',
    'TextBox 23': '"Data Doesn\'t Have a Quality Problem—It Has an Ownership Problem"',
    'TextBox 11': 'OWN',
    'TextBox 13': 'DISCOVER',
    'TextBox 15': 'GOVERN',
    'TextBox 12': 'Domain teams own their data as a product—not IT',
    'TextBox 14': 'Federated catalog makes every asset searchable across domains',
    'TextBox 16': 'Policies travel with the data—access is automatic, not a ticket',
    'TextBox 20': '"Treat data like a product—not a byproduct"',
}

def find_shape(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:
            found = find_shape(shape.shapes, name)
            if found:
                return found
    return None

for name, text in updates.items():
    shape = find_shape(slide.shapes, name)
    if shape and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.runs:
                p.runs[0].text = text
                for run in p.runs[1:]:
                    run.text = ''
                break
        print(f'Updated: {name}')

out = os.path.join(OUTDIR, 'slide21_governance_intro.pptx')
prs.save(out)
print(f'Saved: {out}')
