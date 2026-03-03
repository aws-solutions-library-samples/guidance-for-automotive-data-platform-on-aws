from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def bullet_list(s, l, t, w, h, items, sz=13, color=LIGHT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = False
        p.space_after = Pt(6)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === HEADER ===
circ(slide, 0.5, 0.2, 0.55, PINK, None)
txt(slide, 0.5, 0.24, 0.55, 0.45, "6", 22, DARK, True, PP_ALIGN.CENTER)
txt(slide, 1.2, 0.2, 12, 0.6, "CUSTOMER ENGAGEMENT", 36, ORANGE, True)
txt(slide, 1.2, 0.75, 14, 0.4,
    "Amazon Connect — AI-powered contact center triggered by vehicle events", 22, PINK, True)

# === LEFT — Capabilities ===
lx = 0.6
txt(slide, lx, 1.4, 7, 0.35, "Capabilities", 18, WHITE, True)
rect(slide, lx, 1.85, 7.0, 4.6, CARD, PINK, 1.5)

bullet_list(slide, lx + 0.25, 1.95, 6.5, 4.4, [
    ("Cloud-native contact center that closes the loop between", WHITE),
    ("vehicle events and customer action — proactive, not reactive", LIGHT),
    ("", LIGHT),
    ("▸  Omnichannel — voice, chat, SMS, email, and in-app from a", LIGHT),
    ("   single platform. Customers choose their preferred channel", LIGHT),
    ("", LIGHT),
    ("▸  Event-driven engagement — vehicle DTC triggers automatic", LIGHT),
    ("   outreach: \"We detected low tire pressure, schedule service?\"", LIGHT),
    ("", LIGHT),
    ("▸  Amazon Q in Connect — generative AI agent assist that gives", LIGHT),
    ("   advisors real-time vehicle context, repair history, and", LIGHT),
    ("   recommended next actions during live calls", LIGHT),
    ("", LIGHT),
    ("▸  Contact Lens — real-time sentiment analysis, automatic call", LIGHT),
    ("   summarization, and quality management across every interaction", LIGHT),
    ("", LIGHT),
    ("▸  Outbound campaigns — proactive recall notifications, service", LIGHT),
    ("   reminders, and OTA update confirmations at fleet scale", LIGHT),
], 13, LIGHT)

# === RIGHT — Differentiators ===
rx = 7.9
txt(slide, rx, 1.4, 7, 0.35, "Differentiators", 18, WHITE, True)
rect(slide, rx, 1.85, 7.5, 4.6, CARD, ORANGE, 1.5)

bullet_list(slide, rx + 0.25, 1.95, 7.0, 4.4, [
    ("vs. Genesys / Five9 / legacy CCaaS", WHITE),
    ("No per-seat licensing — pay per minute of use. A 10,000-seat", LIGHT),
    ("contact center costs the same as 10 seats when idle. Scales", LIGHT),
    ("instantly for recall events or seasonal spikes", LIGHT),
    ("", LIGHT),
    ("vs. Salesforce Service Cloud Voice", WHITE),
    ("Connect IS the telephony — not a wrapper. Salesforce integrates", LIGHT),
    ("natively with Connect (Service Cloud Voice) for CRM + contact", LIGHT),
    ("center in one. Best of both worlds, not either/or", LIGHT),
    ("", LIGHT),
    ("Native vehicle data integration", WHITE),
    ("IoT Core events flow directly into Connect via Lambda —", LIGHT),
    ("vehicle context is on the advisor's screen before the customer", LIGHT),
    ("even explains the problem", LIGHT),
    ("", LIGHT),
    ("Pay-as-you-go", WHITE),
    ("Per-minute pricing, no upfront costs, no minimum commitments,", LIGHT),
    ("no long-term contracts. Deploy in minutes, not months", LIGHT),
], 13, LIGHT)

# === BOTTOM ===
py = 6.7
txt(slide, 0.6, py, 4, 0.3, "Automotive Use Cases", 14, MUTED, True)

products = [
    ("Proactive Service", PINK, "Vehicle events trigger outreach before customer calls in"),
    ("Recall Management", RED, "Automated outbound campaigns for safety recalls at scale"),
    ("Roadside Assist", GREEN, "Real-time vehicle location + diagnostics to dispatch teams"),
]

px = 0.6
for name, color, desc in products:
    rect(slide, px, py + 0.35, 4.8, 0.7, CARD, color, 1)
    txt(slide, px + 0.2, py + 0.38, 2.5, 0.3, name, 14, WHITE, True)
    txt(slide, px + 0.2, py + 0.65, 4.4, 0.3, desc, 11, LIGHT)
    px += 5.0

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide15_connect.pptx")
prs.save(out)
print(f"Saved: {out}")
