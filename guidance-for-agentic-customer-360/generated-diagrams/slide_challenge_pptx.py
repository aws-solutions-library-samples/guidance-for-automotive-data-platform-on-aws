from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
RED = RGBColor(248, 113, 113)
YELLOW = RGBColor(251, 191, 36)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(16, 20, 36)

def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    if align: p.alignment = align

def mtxt(slide, l, t, w, h, lines):
    """Multi-format text: list of (text, size, color, bold) tuples, each on new line."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, sz, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(4)

def rect(slide, l, t, w, h, fill=CARD, border=None, border_width=1.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(border_width)
    else: s.line.fill.background()

# Title
txt(slide, 0.6, 0.3, 12, 0.7, "WHAT'S HOLDING YOU AT GEN 2?", 36, ORANGE, True)
txt(slide, 0.6, 0.85, 14, 0.5, "The barriers between cloud-connected and cloud-native & AI-enabled", 20, LIGHT)

# Gen 2 → Gen 3 context bar
rect(slide, 0.6, 1.5, 7, 0.5, CARD, YELLOW)
txt(slide, 0.9, 1.55, 6, 0.4, "Gen 2: Cloud-Connected  (where most OEMs are today)", 14, YELLOW, True)
rect(slide, 8.4, 1.5, 7, 0.5, CARD, GREEN)
txt(slide, 8.7, 1.55, 6, 0.4, "Gen 3: Cloud-Native & AI-Enabled  (where you need to be)", 14, GREEN, True)
# Arrow between
txt(slide, 7.65, 1.48, 0.8, 0.5, "→", 24, MUTED, True)

# Four challenge cards — 2x2 grid, large and impactful
challenges = [
    {
        "icon": "📈",
        "title": "Scaling Ceiling",
        "stat": "10x",
        "stat_desc": "fleet growth expected",
        "desc": "Self-managed infrastructure hits hard limits as fleets grow from thousands to millions of vehicles. Every new region, every new model year compounds the problem.",
        "gen2": "Manual capacity planning",
        "gen3": "Auto-scaling managed services",
        "color": SKY,
    },
    {
        "icon": "🔧",
        "title": "Ops Burden",
        "stat": "70%",
        "stat_desc": "of eng time on infra",
        "desc": "Engineering teams spend the majority of their time patching, scaling, and maintaining infrastructure instead of building features that differentiate your vehicles.",
        "gen2": "Teams manage servers & pipelines",
        "gen3": "Zero-ops fully managed services",
        "color": RED,
    },
    {
        "icon": "💰",
        "title": "Cost Doesn't Scale",
        "stat": "Linear",
        "stat_desc": "cost growth with fleet",
        "desc": "On-prem and self-managed cloud costs grow linearly with fleet size. Every additional vehicle adds proportional infrastructure cost with no economies of scale.",
        "gen2": "Pay for peak capacity 24/7",
        "gen3": "Pay-as-you-go, scale to zero",
        "color": YELLOW,
    },
    {
        "icon": "⚡",
        "title": "Innovation Gap",
        "stat": "Weeks",
        "stat_desc": "vs. quarters to ship",
        "desc": "Competitors on managed services ship new features in days. You're still waiting on infrastructure provisioning, capacity reviews, and deployment windows.",
        "gen2": "Quarterly release cycles",
        "gen3": "Continuous delivery in days",
        "color": GREEN,
    },
]

cw = 7.0
ch = 2.8
cgap = 0.4
cx_start = 0.6
cy_start = 2.3

for i, c in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = cx_start + col * (cw + cgap)
    y = cy_start + row * (ch + 0.3)

    # Card
    rect(slide, x, y, cw, ch, CARD, c["color"], 1)

    # Icon + Title
    txt(slide, x+0.2, y+0.15, 0.5, 0.4, c["icon"], 22, WHITE)
    txt(slide, x+0.6, y+0.15, 4, 0.4, c["title"], 22, WHITE, True)

    # Big stat
    txt(slide, x+cw-2.2, y+0.1, 2, 0.5, c["stat"], 32, c["color"], True)
    txt(slide, x+cw-2.2, y+0.5, 2, 0.3, c["stat_desc"], 11, MUTED)

    # Description
    txt(slide, x+0.3, y+0.75, cw-0.6, 0.9, c["desc"], 13, LIGHT)

    # Gen 2 vs Gen 3 comparison
    rect(slide, x+0.2, y+ch-0.8, cw/2-0.3, 0.6, DARK, YELLOW, 0.75)
    txt(slide, x+0.35, y+ch-0.78, cw/2-0.5, 0.2, "Gen 2", 10, YELLOW, True)
    txt(slide, x+0.35, y+ch-0.55, cw/2-0.5, 0.3, c["gen2"], 11, MUTED)

    rect(slide, x+cw/2+0.1, y+ch-0.8, cw/2-0.3, 0.6, DARK, GREEN, 0.75)
    txt(slide, x+cw/2+0.25, y+ch-0.78, cw/2-0.5, 0.2, "Gen 3", 10, GREEN, True)
    txt(slide, x+cw/2+0.25, y+ch-0.55, cw/2-0.5, 0.3, c["gen3"], 11, WHITE)

# Bottom CTA
rect(slide, 0.6, 8.0, 14.8, 0.45, CARD, ORANGE)
txt(slide, 0.9, 8.05, 10, 0.35, "The path from Gen 2 to Gen 3 isn't a rewrite — it's a migration to managed services.", 16, LIGHT)
txt(slide, 11.5, 8.05, 4, 0.35, "AWS makes it possible.", 16, ORANGE, True)

# Footer
txt(slide, 0.6, 8.65, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide_challenge.pptx")
prs.save(out)
print(f"Saved: {out}")
