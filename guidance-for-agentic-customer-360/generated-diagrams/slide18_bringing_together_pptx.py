from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = '/Users/givenand/Downloads/CVx-FCD-2026-slide18.pptx'
OUTDIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation(TEMPLATE)
slide = prs.slides[0]

# Remove all shapes except the ones we want to keep
keep_names = {
    "Freeform 2", "Group 3", "AutoShape 5", "Group 6",
    "TextBox 64", "TextBox 65", "TextBox 66", "TextBox 67"
}

spTree = slide.shapes._spTree
for el in [s._element for s in slide.shapes if s.name not in keep_names]:
    spTree.remove(el)

# Update header
for shape in slide.shapes:
    if shape.name == "TextBox 65":
        shape.text_frame.paragraphs[0].text = "BRINGING IT ALL TOGETHER"
    elif shape.name == "TextBox 66":
        shape.text_frame.paragraphs[0].text = "From fragmented silos to a unified automotive data platform"
    elif shape.name == "TextBox 64":
        shape.text_frame.paragraphs[0].text = "4"

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(l, t, sz, fill=CARD, border=None, bw=2):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow_line(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TOP ROW — 5 data silos ===
silos = [
    ("CRM", SKY),
    ("DMS", GREEN),
    ("Telemetry", YELLOW),
    ("Service", PURPLE),
    ("Contact Center", PINK),
]

sw = 2.8
sgap = 0.55
sx = 1.4
sy = 2.0

for i, (name, color) in enumerate(silos):
    x = sx + i * (sw + sgap)
    rect(x, sy, sw, 0.9, CARD, color, 1.5)
    txt(x + 0.1, sy + 0.15, sw - 0.2, 0.6, name, 14, WHITE, True, PP_ALIGN.CENTER)

# Arrows down from silos
for i in range(5):
    x = sx + i * (sw + sgap) + sw / 2
    arrow_line(x, sy + 0.9, x, sy + 1.4, MUTED, 1.5)

# === INGESTION LAYER ===
ing_y = sy + 1.4
rect(1.0, ing_y, 18.0, 0.7, CARD, ORANGE, 1.5)
txt(1.3, ing_y + 0.08, 3.5, 0.25, "Data Ingestion", 15, ORANGE, True)
txt(5.0, ing_y + 0.08, 13.5, 0.25,
    "IoT Core (MQTT)  ·  Amazon MSK (Kafka)  ·  Kinesis  ·  API Gateway  ·  Glue ETL", 12, LIGHT)
txt(1.3, ing_y + 0.38, 17, 0.25,
    "Real-time streaming, near-real-time APIs, and batch processing into a unified data lake", 10, MUTED)

# Arrow down
arrow_line(10.0, ing_y + 0.7, 10.0, ing_y + 1.15, ORANGE, 2)

# === ADP PLATFORM BOX ===
plat_y = ing_y + 1.15
plat_h = 3.6
rect(1.0, plat_y, 18.0, plat_h, CARD, ORANGE, 2)

txt(1.3, plat_y + 0.1, 17, 0.35, "Automotive Data Platform on AWS", 18, ORANGE, True, PP_ALIGN.CENTER)

# 4 pillar cards — with icon placeholder space at top
pillars = [
    ("Customer 360\nAnalytics", "Bedrock Agent · QuickSight\nAurora pgvector · Athena", SKY, "← Add Bedrock icon"),
    ("Predictive\nMaintenance", "SageMaker · IoT FleetWise\nGlue ETL · Step Functions", GREEN, "← Add SageMaker icon"),
    ("Data Mesh\nFoundation", "SageMaker Unified Studio\nDataZone · Glue Catalog", PURPLE, "← Add DataZone icon"),
    ("Data\nGovernance", "Lake Formation · Macie\nCloudTrail · IAM", PINK, "← Add Lake Formation icon"),
]

pw = 3.9
pgap = 0.3
ppx = 1.5
ppy = plat_y + 0.55

for name, services, color, icon_hint in pillars:
    rect(ppx, ppy, pw, 2.8, DARK, color, 1.5)

    # Icon placeholder circle — swap for AWS service icon in Canva
    icon_cx = ppx + pw / 2 - 0.4
    circ(icon_cx, ppy + 0.15, 0.8, RGBColor(25, 35, 60), color, 1.5)
    txt(icon_cx, ppy + 0.25, 0.8, 0.6, "?", 24, color, True, PP_ALIGN.CENTER)

    # Pillar name
    txt(ppx + 0.15, ppy + 1.05, pw - 0.3, 0.6, name, 14, WHITE, True, PP_ALIGN.CENTER)

    # Services
    txt(ppx + 0.15, ppy + 1.7, pw - 0.3, 0.6, services, 10, LIGHT, False, PP_ALIGN.CENTER)

    # Hint text (remove in final)
    txt(ppx + 0.15, ppy + 2.4, pw - 0.3, 0.3, icon_hint, 8, MUTED, False, PP_ALIGN.CENTER)

    ppx += pw + pgap

# Arrow down from platform
arrow_line(10.0, plat_y + plat_h, 10.0, plat_y + plat_h + 0.35, ORANGE, 2)

# === OUTCOMES ROW — icons + labels, no boxes ===
oy = plat_y + plat_h + 0.35
outcomes = [
    ("Proactive\nRetention", "Identify at-risk customers\nbefore they churn", SKY),
    ("Predictive\nService", "7-14 day advance warning\nof failures", GREEN),
    ("AI-Powered\nInsights", "Natural language queries\nacross all data", PURPLE),
    ("Compliant\nData Sharing", "EU Data Act + GDPR\nbuilt in", PINK),
]

ow = 3.9
ogap = 0.3
ox = 1.5

for name, desc, color in outcomes:
    # Icon placeholder circle — swap for outcome icon in Canva
    icon_cx = ox + ow / 2 - 0.35
    circ(icon_cx, oy, 0.7, DARK, color, 1.5)
    txt(icon_cx, oy + 0.08, 0.7, 0.5, "?", 22, color, True, PP_ALIGN.CENTER)

    # Outcome name
    txt(ox, oy + 0.8, ow, 0.5, name, 12, WHITE, True, PP_ALIGN.CENTER)

    # Description
    txt(ox, oy + 1.3, ow, 0.5, desc, 9, MUTED, False, PP_ALIGN.CENTER)

    ox += ow + ogap

out = os.path.join(OUTDIR, "slide18_bringing_together.pptx")
prs.save(out)
print(f"Saved: {out}")
