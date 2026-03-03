from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=11, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.6, "DATA GOVERNANCE & COMPLIANCE DEEP DIVE", 36, ORANGE, True)
txt(0.6, 0.75, 14, 0.4,
    "Multi-region data governance with fine-grained access control and automated PII detection", 18, LIGHT)

# === TOP: Three-Account Architecture ===
txt(0.6, 1.35, 8, 0.3, "MULTI-REGION ACCOUNT ARCHITECTURE", 13, ORANGE, True)

# Central Governance Account
rect(4.5, 1.7, 7.0, 1.6, CARD, ORANGE, 2)
txt(4.8, 1.75, 6.5, 0.25, "Central Governance Account", 14, ORANGE, True)
mtxt(4.8, 2.05, 6.5, 1.2, [
    ("Lake Formation — fine-grained access policies", WHITE, False),
    ("DataZone — federated data catalog", LIGHT, False),
    ("CloudTrail — audit trail for all data access", LIGHT, False),
    ("RAM — cross-account resource sharing", LIGHT, False),
], 10)

# Arrows down to producer and consumer
arrow(6.5, 3.3, 3.5, 3.7, ORANGE, 1.5)
arrow(9.5, 3.3, 12.5, 3.7, ORANGE, 1.5)

# EU Producer Account
rect(0.6, 3.7, 6.5, 2.0, CARD, SKY, 1.5)
txt(0.9, 3.75, 6, 0.25, "EU Producer Account (eu-west-1)", 13, SKY, True)
mtxt(0.9, 4.05, 6, 1.6, [
    ("Data residency: EU customer & service data", WHITE, False),
    ("", LIGHT, False),
    ("S3 — EU data lake (encrypted, versioned)", LIGHT, False),
    ("Glue — ETL & catalog for EU datasets", LIGHT, False),
    ("Macie — automated PII scanning", PINK, False),
    ("  → Names, emails, VINs, addresses detected", MUTED, False),
    ("  → Auto-classification & remediation alerts", MUTED, False),
], 10)

# Global Consumer Account
rect(8.9, 3.7, 6.5, 2.0, CARD, GREEN, 1.5)
txt(9.2, 3.75, 6, 0.25, "Global Consumer Account (us-east-1)", 13, GREEN, True)
mtxt(9.2, 4.05, 6, 1.6, [
    ("Aggregated & anonymized views only", WHITE, False),
    ("", LIGHT, False),
    ("Athena — cross-region federated queries", LIGHT, False),
    ("QuickSight — global executive dashboards", LIGHT, False),
    ("Bedrock — AI insights on aggregated data", LIGHT, False),
    ("  → No PII exposed in global views", MUTED, False),
    ("  → Row/column filtering via Lake Formation", MUTED, False),
], 10)

# === MIDDLE: Access Control Scenarios ===
txt(0.6, 6.0, 8, 0.3, "ACCESS CONTROL IN ACTION", 13, YELLOW, True)

# Scenario 1: German service rep
rect(0.6, 6.35, 7.3, 1.1, CARD, SKY, 1.5)
txt(0.9, 6.4, 6.8, 0.25, "Persona: German Service Representative", 12, SKY, True)
mtxt(0.9, 6.7, 6.8, 0.7, [
    ("Query: \"Show me customer complaints for Berlin dealers\"", WHITE, False),
    ("→ Lake Formation row filter: region = 'EU-DE'", GREEN, False),
    ("→ Sees: EU customer data, service history, VIN details", LIGHT, False),
    ("→ Cannot see: US data, global aggregates, financial metrics", RED, False),
], 10)

# Scenario 2: Global exec
rect(8.3, 6.35, 7.1, 1.1, CARD, GREEN, 1.5)
txt(8.6, 6.4, 6.5, 0.25, "Persona: Global VP of Customer Experience", 12, GREEN, True)
mtxt(8.6, 6.7, 6.5, 0.7, [
    ("Query: \"Compare NPS across all regions\"", WHITE, False),
    ("→ Lake Formation column filter: no PII columns", YELLOW, False),
    ("→ Sees: Aggregated NPS, trends, regional benchmarks", LIGHT, False),
    ("→ Cannot see: Individual names, emails, addresses", RED, False),
], 10)

# === BOTTOM: Compliance Bar ===
rect(0.6, 7.7, 14.8, 0.55, CARD, ORANGE, 1)
compliance = [
    ("EU Data Act", "vehicle data portability", 0.9),
    ("GDPR", "right to erasure & minimization", 3.2),
    ("Macie", "automated PII detection", 5.8),
    ("CloudTrail", "full audit trail", 8.2),
    ("Lake Formation", "row/column ACLs", 10.5),
    ("Multi-region", "data residency enforced", 13.0),
]
for val, label, x in compliance:
    txt(x, 7.75, 2.2, 0.2, val, 12, ORANGE, True)
    txt(x, 7.97, 2.2, 0.2, label, 9, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide21_governance_deep_dive.pptx")
prs.save(out)
print(f"Saved: {out}")
