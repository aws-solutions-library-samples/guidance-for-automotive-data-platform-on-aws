from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def bullet_list(s, l, t, w, h, items, sz=13, color=LIGHT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = False
        p.space_after = Pt(6)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def pill(s, l, t, text, color, sz=11):
    pw = len(text) * 0.085 + 0.35
    ph = 0.3
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(pw), Inches(ph))
    dark = RGBColor(color[0] // 6, color[1] // 6, color[2] // 6)
    sh.fill.solid(); sh.fill.fore_color.rgb = dark
    sh.line.color.rgb = color; sh.line.width = Pt(1)
    tb = sh.text_frame; tb.word_wrap = False
    p = tb.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz); p.font.color.rgb = color
    return pw

# === HEADER — component number + title ===
circ(slide, 0.5, 0.2, 0.55, SKY, None)
txt(slide, 0.5, 0.24, 0.55, 0.45, "1", 22, DARK, True, PP_ALIGN.CENTER)
txt(slide, 1.2, 0.2, 12, 0.6, "CELLULAR CONNECTIVITY", 36, ORANGE, True)
txt(slide, 1.2, 0.75, 12, 0.4, "Amazon Connectivity Services (ACS)", 22, SKY, True)

# === LEFT COLUMN — What it does ===
lx = 0.6
txt(slide, lx, 1.4, 7, 0.35, "What ACS Does", 18, WHITE, True)

rect(slide, lx, 1.85, 7.0, 4.6, CARD, SKY, 1.5)

bullet_list(slide, lx + 0.25, 1.95, 6.5, 4.4, [
    ("Managed connectivity orchestration for connected vehicles", WHITE),
    ("at scale — from cellular to Wi-Fi to satellite", LIGHT),
    ("", LIGHT),
    ("▸  Multi-MNO profile management with sub-60-second", LIGHT),
    ("   automated switching via cloud-based eSIM (SGP.31/32)", LIGHT),
    ("", LIGHT),
    ("▸  Network Coverage & Quality SDK — real-time signal", LIGHT),
    ("   monitoring, billing validation, and anomaly detection", LIGHT),
    ("", LIGHT),
    ("▸  Wi-Fi offloading across 50M+ access points (US),", LIGHT),
    ("   reducing data costs 50-70% for non-latency-sensitive traffic", LIGHT),
    ("", LIGHT),
    ("▸  Multi-connectivity ready: cellular, Wi-Fi, Amazon", LIGHT),
    ("   Sidewalk, and Amazon LEO satellite (Kuiper)", LIGHT),
    ("", LIGHT),
    ("▸  Single pane of glass — unified management across", LIGHT),
    ("   all carriers, networks, and connectivity types", LIGHT),
], 13, LIGHT)

# === RIGHT COLUMN — Why ACS ===
rx = 7.9
txt(slide, rx, 1.4, 7, 0.35, "Why ACS", 18, WHITE, True)

rect(slide, rx, 1.85, 7.5, 4.6, CARD, ORANGE, 1.5)

bullet_list(slide, rx + 0.25, 1.95, 7.0, 4.4, [
    ("Proven at Amazon scale", WHITE),
    ("Manages connectivity for Amazon Last Mile fleet, Zoox", LIGHT),
    ("autonomous vehicles, and 40+ Amazon business customers", LIGHT),
    ("", LIGHT),
    ("Eliminates aggregator margins", WHITE),
    ("Direct tier 1 MNO partnerships remove 15-30% aggregator", LIGHT),
    ("markup — achieving automotive market rates ($0.60-1.20/GB)", LIGHT),
    ("", LIGHT),
    ("Creates MNO negotiating leverage", WHITE),
    ("SGP.31/32 eSIM enables dynamic traffic reallocation —", LIGHT),
    ("MNOs know you can shift traffic to competitors in minutes", LIGHT),
    ("", LIGHT),
    ("Transforms connectivity from cost center to advantage", WHITE),
    ("Data consumption growing 3x+ every 2 years (5 GB → 207 GB", LIGHT),
    ("per vehicle by 2029) — infrastructure must scale with it", LIGHT),
], 13, LIGHT)

# === BOTTOM — Key Products ===
py = 6.7
txt(slide, 0.6, py, 3, 0.3, "Key Products", 14, MUTED, True)

products = [
    ("Conekt Platform", SKY, "Connectivity orchestration & MNO management"),
    ("Network Quality SDK", GREEN, "Device-side monitoring & billing validation"),
    ("SGP.31/32 eSIM", PURPLE, "Cloud-based multi-profile eSIM infrastructure"),
]

px = 0.6
for name, color, desc in products:
    rect(slide, px, py + 0.35, 4.8, 0.7, CARD, color, 1)
    txt(slide, px + 0.2, py + 0.38, 2.5, 0.3, name, 14, WHITE, True)
    txt(slide, px + 0.2, py + 0.65, 4.4, 0.3, desc, 11, LIGHT)
    px += 5.0

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide10_acs.pptx")
prs.save(out)
print(f"Saved: {out}")
