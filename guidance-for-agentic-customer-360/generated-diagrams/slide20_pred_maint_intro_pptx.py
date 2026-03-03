from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, copy

TEMPLATE = '/Users/givenand/Downloads/CVx-FCD-2026.pptx.pptx'
OUTDIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation(TEMPLATE)
slide = prs.slides[0]

# Map existing shape names to new content
updates = {
    # Header
    'TextBox 10': 'Predictive Maintenance: The solution',
    # Subtitle
    'TextBox 23': '"Maintenance Schedules Are Reactive—Tire Failures Shouldn\'t Be"',
    # Three pillars - keywords
    'TextBox 11': 'DETECT',
    'TextBox 13': 'VALIDATE',
    'TextBox 15': 'PREVENT',
    # Three pillars - descriptions
    'TextBox 12': 'ML models scan tire telemetry for anomalies 24/7',
    'TextBox 14': 'Dual-path confirmation: ML scoring + physics-based leak analysis',
    'TextBox 16': 'Automated alerts trigger service scheduling before failure',
    # Tagline
    'TextBox 20': '"Predict failures 7-14 days out—not after the breakdown"',
}

def update_text(shape, new_text):
    """Update text while preserving all formatting."""
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.runs:
                # Keep first run's formatting, clear the rest
                p.runs[0].text = new_text
                for run in p.runs[1:]:
                    run.text = ''
                return

def find_shape_recursive(shapes, name):
    """Find shape by name, including inside groups."""
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:  # GROUP
            found = find_shape_recursive(shape.shapes, name)
            if found:
                return found
    return None

for name, text in updates.items():
    shape = find_shape_recursive(slide.shapes, name)
    if shape:
        update_text(shape, text)
        print(f'Updated: {name} → "{text[:50]}..."')
    else:
        print(f'NOT FOUND: {name}')

# Update slide number if present
num_shape = find_shape_recursive(slide.shapes, 'TextBox 64')
if num_shape:
    update_text(num_shape, '5')

out = os.path.join(OUTDIR, 'slide20_pred_maint_intro.pptx')
prs.save(out)
print(f'\nSaved: {out}')
