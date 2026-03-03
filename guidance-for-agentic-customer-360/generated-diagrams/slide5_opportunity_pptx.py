from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(16, 20, 36)

def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(slide, l, t, w, h, fill=CARD, border=None, bw=1.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else: s.line.fill.background()

# === TITLE ===
txt(slide, 0.6, 0.3, 12, 0.7, "THE CONNECTED VEHICLE OPPORTUNITY", 36, ORANGE, True)
txt(slide, 0.6, 0.8, 12, 0.5, "Why connected vehicle platforms matter more than ever", 20, LIGHT)

# === TOP THIRD: Stat cards ===
stats = [
    ("95%", "of new vehicles\nconnected by 2030", "Source: McKinsey", SKY),
    ("470M", "connected cars\non the road by 2025", "Source: Statista", GREEN),
    ("$100B+", "connected car services\nmarket by 2030", "Source: McKinsey", ORANGE),
    ("25 GB", "of data per vehicle\nper hour", "Source: Intel / Frost & Sullivan", YELLOW),
]
sw = 3.5; sgap = 0.27; sx = 0.6
for i, (num, desc, source, color) in enumerate(stats):
    x = sx + i * (sw + sgap)
    rect(slide, x, 1.5, sw, 1.5, CARD, color)
    txt(slide, x+0.2, 1.55, 3, 0.6, num, 34, color, True)
    for j, line in enumerate(desc.split("\n")):
        txt(slide, x+0.2, 2.1+j*0.25, 3, 0.3, line, 14, WHITE)
    txt(slide, x+0.2, 2.7, 3, 0.3, source, 10, MUTED)

# === MIDDLE THIRD: Gen 1 / 2 / 3 ===
ew = 4.7; egap = 0.27
gens = [
    ("Gen 1: Basic Telematics", SKY,
     ["One-way data uploads", "Reactive diagnostics", "Siloed on-prem systems", "Manual fleet monitoring"]),
    ("Gen 2: Cloud-Connected", YELLOW,
     ["Real-time bi-directional streaming", "OTA software updates", "Cloud-hosted but self-managed", "Custom-built data pipelines"]),
    ("Gen 3: Cloud-Native & AI-Enabled", GREEN,
     ["Fully managed services (zero ops)", "Real-time stream processing at scale", "AI/ML-driven predictions & automation", "Agentic AI across vehicle lifecycle"]),
]
for i, (label, color, items) in enumerate(gens):
    x = 0.6 + i * (ew + egap)
    rect(slide, x, 3.3, ew, 2.3, CARD, color)
    txt(slide, x+0.2, 3.4, 4.3, 0.4, label, 16, color, True)
    for j, item in enumerate(items):
        txt(slide, x+0.2, 3.85+j*0.38, 4.3, 0.35, f"→  {item}", 13, LIGHT)

# === BOTTOM THIRD: Challenge — Gen 2 → Gen 3 barriers ===
rect(slide, 0.6, 5.9, 14.8, 1.5, CARD, SKY, 1)
txt(slide, 0.9, 5.98, 8, 0.4, "WHAT'S HOLDING YOU AT GEN 2?", 18, SKY, True)

challenges = [
    ("📈  Scaling Ceiling", "10x fleet growth — self-managed infra hits limits at millions of vehicles"),
    ("🔧  Ops Burden", "70% of eng time on infrastructure — not building differentiating features"),
    ("💰  Cost Doesn't Scale", "Linear cost growth with fleet — no economies of scale on self-managed"),
    ("⚡  Innovation Gap", "Competitors ship in weeks on managed services — you're waiting on infra"),
]
cw2 = 7.0
for i, (title, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = 0.9 + col * cw2
    y = 6.45 + row * 0.42
    txt(slide, x, y, 2.5, 0.35, title, 13, WHITE, True)
    txt(slide, x+2.5, y, 4.3, 0.35, desc, 12, LIGHT)

# === BOTTOM HOOK ===
rect(slide, 0.6, 7.7, 14.8, 0.4, CARD)
txt(slide, 0.9, 7.73, 7, 0.35, "The path from Gen 2 to Gen 3 isn't a rewrite —", 15, LIGHT)
txt(slide, 7.2, 7.73, 5, 0.35, "it's a migration to managed services.", 15, ORANGE, True)

# Footer
txt(slide, 0.6, 8.6, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide5_opportunity.pptx")
prs.save(out)
print(f"Saved: {out}")
