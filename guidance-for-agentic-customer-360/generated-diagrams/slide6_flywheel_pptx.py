from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import math, os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(slide, l, t, w, h, fill=CARD, border=None, bw=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else: s.line.fill.background()
    return s

def circle(slide, l, t, size, fill=CARD, border=None, bw=2):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(size), Inches(size))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else: s.line.fill.background()
    return s

def arrow(slide, x1, y1, x2, y2, color=SKY, width=2):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

# === TITLE ===
txt(slide, 0.6, 0.25, 10, 0.6, "CONNECTED VEHICLE PLATFORM FLYWHEEL", 36, ORANGE, True)
txt(slide, 0.6, 0.8, 12, 0.4, "Six managed AWS services powering the continuous vehicle data loop", 20, LIGHT)

# === FLYWHEEL — 6 nodes in a circle ===
# Center of flywheel
cx, cy = 5.5, 4.8
radius = 2.8
node_size = 1.6

# 6 services positioned around the circle (starting from top, clockwise)
services = [
    ("ACS\nCellular", "Managed cellular\nconnectivity", SKY, "1"),
    ("IoT Core\nConnectivity", "Billions of messages\nauto-scaling", GREEN, "2"),
    ("Amazon MSK\nStreaming", "Managed Kafka\nunlimited throughput", YELLOW, "3"),
    ("Managed Flink\nProcessing", "Real-time analytics\nsub-second latency", ORANGE, "4"),
    ("S3 + Athena\nAnalytics", "Serverless storage\nSQL on petabytes", PURPLE, "5"),
    ("Amazon Connect\nEngagement", "AI-powered IVR\nomnichannel", PINK, "6"),
]

# Calculate positions (top = -90 degrees, clockwise)
positions = []
for i in range(6):
    angle = math.radians(-90 + i * 60)
    x = cx + radius * math.cos(angle) - node_size/2
    y = cy + radius * math.sin(angle) - node_size/2
    positions.append((x, y))

# Draw connecting arrows between nodes (clockwise)
for i in range(6):
    next_i = (i + 1) % 6
    x1 = positions[i][0] + node_size/2
    y1 = positions[i][1] + node_size/2
    x2 = positions[next_i][0] + node_size/2
    y2 = positions[next_i][1] + node_size/2
    # Shorten arrows to not overlap circles
    dx = x2 - x1; dy = y2 - y1
    dist = math.sqrt(dx*dx + dy*dy)
    offset = node_size * 0.55
    sx = x1 + dx/dist * offset
    sy = y1 + dy/dist * offset
    ex = x2 - dx/dist * offset
    ey = y2 - dy/dist * offset
    color = services[i][2]
    arrow(slide, sx, sy, ex, ey, color, 2.5)

# Draw nodes
for i, (name, desc, color, num) in enumerate(services):
    x, y = positions[i]

    # Outer circle
    circle(slide, x, y, node_size, CARD, color, 2.5)

    # Number badge
    circle(slide, x - 0.15, y - 0.15, 0.4, color, None)
    txt(slide, x - 0.15, y - 0.12, 0.4, 0.35, num, 14, DARK, True, PP_ALIGN.CENTER)

    # Service name
    txt(slide, x + 0.1, y + 0.25, node_size - 0.2, 0.7, name, 12, WHITE, True, PP_ALIGN.CENTER)

    # Description
    txt(slide, x + 0.1, y + 0.85, node_size - 0.2, 0.6, desc, 9, MUTED, False, PP_ALIGN.CENTER)

# === CENTER — Flywheel label ===
circle(slide, cx - 0.9, cy - 0.9, 1.8, RGBColor(16, 22, 40), ORANGE, 2)
txt(slide, cx - 0.85, cy - 0.55, 1.7, 0.4, "Vehicle", 16, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.85, cy - 0.2, 1.7, 0.4, "Data", 16, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.85, cy + 0.15, 1.7, 0.4, "Flywheel", 16, ORANGE, True, PP_ALIGN.CENTER)

# === RIGHT SIDE — Flow description ===
rx = 10.0
rect(slide, rx, 1.5, 5.4, 6.8, CARD)

stages = [
    ("Ingest & Connect", "Vehicle signals flow through managed cellular (ACS) and IoT Core — no brokers to manage", SKY, "1 → 2"),
    ("Stream & Process", "Amazon MSK ingests events, Managed Flink processes in real-time — zero ops Kafka & Flink", YELLOW, "3 → 4"),
    ("Analyze & Act", "S3 + Athena for serverless analytics, Amazon Connect for customer engagement & alerts", PURPLE, "5 → 6"),
    ("Loop Back", "Insights drive OTA updates, proactive service, and personalized experiences back to the vehicle", ORANGE, "6 → 1"),
]

sy = 1.7
for i, (title, desc, color, nums) in enumerate(stages):
    y = sy + i * 1.6

    # Stage number pill
    rect(slide, rx + 0.2, y, 0.8, 0.35, DARK, color, 1)
    txt(slide, rx + 0.25, y + 0.03, 0.7, 0.3, nums, 11, color, True, PP_ALIGN.CENTER)

    # Title
    txt(slide, rx + 1.2, y, 4, 0.35, title, 16, WHITE, True)

    # Description
    txt(slide, rx + 1.2, y + 0.4, 4, 0.8, desc, 12, LIGHT)

# === OBSERVABILITY BAR ===
rect(slide, rx + 0.2, 8.0, 5.0, 0.35, DARK, RED, 1)
txt(slide, rx + 0.4, 8.03, 4.5, 0.3, "☁  Amazon CloudWatch — observability across all layers", 12, RED)

# Footer
txt(slide, 0.6, 8.6, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide6_flywheel.pptx")
prs.save(out)
print(f"Saved: {out}")
