from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = '/Users/givenand/Downloads/2026_aws_powerpoint_template_v1_07162eba (1).pptx'
OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

prs = Presentation(TEMPLATE)
# Clear existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Use Title Only Gradient 1 (dark background with title)
slide = prs.slides.add_slide(prs.slide_layouts[13])

# Title
slide.placeholders[0].text = "Automotive Data Platform — Architecture Overview"
for p in slide.placeholders[0].text_frame.paragraphs:
    p.font.color.rgb = WHITE

def txt(s, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# Subtitle
txt(slide, 0.5, 0.85, 12, 0.3,
    "Four integrated solutions on a six-layer serverless architecture", 13, MUTED)

# === Four pillar cards ===
pillars = [
    ("Customer 360\nAnalytics", "Bedrock Agent\nQuickSight", SKY),
    ("Predictive\nMaintenance", "SageMaker\nIoT Core", GREEN),
    ("Data Mesh\nFoundation", "Unified Studio\nDataZone", PURPLE),
    ("Data\nGovernance", "Lake Formation\nMacie", PINK),
]

pw = 2.9
pgap = 0.2
px = 0.5
py = 1.3

for name, services, color in pillars:
    rect(slide, px, py, pw, 1.4, CARD, color, 2)
    txt(slide, px + 0.15, py + 0.1, pw - 0.3, 0.55, name, 13, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, px + 0.15, py + 0.75, pw - 0.3, 0.5, services, 10, color, False, PP_ALIGN.CENTER)
    px += pw + pgap

# === Six-layer architecture stack ===
layers = [
    ("Data Sources", "CRM · DMS · Telemetry · Service · Contact Center · Third-Party", ORANGE),
    ("Data Ingestion", "IoT Core (MQTT) · Amazon MSK (Kafka) · Kinesis · API Gateway", SKY),
    ("Storage & Governance", "S3 Data Lake · Glue Catalog · Lake Formation · DataZone", GREEN),
    ("Processing", "Managed Flink (real-time) · Glue ETL (batch) · SageMaker (ML)", PURPLE),
    ("Analytics & AI", "Athena (SQL) · QuickSight (dashboards) · Bedrock (AI agents)", YELLOW),
    ("Action", "Amazon Connect (proactive outreach) · EventBridge (alerts) · SNS", PINK),
]

ly = 2.95
lh = 0.5
lgap = 0.08

for name, services, color in layers:
    # Color bar
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(ly), Inches(0.08), Inches(lh))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

    # Row
    rect(slide, 0.65, ly, 12.15, lh, CARD)

    # Layer name
    txt(slide, 0.8, ly + 0.05, 2.2, 0.35, name, 12, WHITE, True)

    # Services
    txt(slide, 3.1, ly + 0.05, 9.5, 0.35, services, 11, LIGHT)

    ly += lh + lgap

# Stats bar
sy = ly + 0.15
rect(slide, 0.5, sy, 12.3, 0.45, CARD, ORANGE, 1)
stats = [
    ("500K customers", 0.7),
    ("1.4M interactions", 3.0),
    ("900K service records", 5.3),
    ("11 datasets", 7.8),
    ("~$114/mo", 9.5),
    ("MIT-0 License", 11.0),
]
for label, x in stats:
    txt(slide, x, sy + 0.07, 2, 0.3, label, 11, ORANGE, True)

prs.save(os.path.join(OUTDIR, "slide18_adp_arch.pptx"))
print("Saved: slide18_adp_arch.pptx (dark gradient template)")
