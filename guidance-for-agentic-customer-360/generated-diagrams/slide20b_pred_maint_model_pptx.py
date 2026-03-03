from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — MODEL & INFERENCE", 34, ORANGE, True)
txt(0.6, 0.65, 14, 0.35,
    "SageMaker Random Cut Forest training, batch inference, real-time API, and anomaly thresholding", 17, LIGHT)

# =====================================================================
# TOP LEFT: TRAINING DATASET & FEATURES
# =====================================================================
txt(0.6, 1.2, 7, 0.25, "TRAINING DATASET & FEATURE ENGINEERING", 13, ORANGE, True)

rect(0.6, 1.5, 7.3, 2.3, CARD, PURPLE, 1.5)

mtxt(0.9, 1.55, 3.3, 2.2, [
    ("Raw Telemetry (from Redshift):", WHITE, True),
    ("• pressure — tire PSI reading", LIGHT, False),
    ("• temperature — tire temp (°F)", LIGHT, False),
    ("• vehicle_id — unique vehicle", LIGHT, False),
    ("• tire_id — position (FL/FR/RL/RR)", LIGHT, False),
    ("• timestamp — event time", LIGHT, False),
    ("", LIGHT, False),
    ("Root ETL Output:", WHITE, True),
    ("Cleaned CSV, partitioned by date/hour", MUTED, False),
], 10)

mtxt(4.2, 1.55, 3.5, 2.2, [
    ("ML ETL Engineered Features:", WHITE, True),
    ("• delta_pressure — lag(pressure) diff", LIGHT, False),
    ("• delta_temp — lag(temperature) diff", LIGHT, False),
    ("", LIGHT, False),
    ("Normalization (z-score):", WHITE, True),
    ("• (value − mean) / std per feature", LIGHT, False),
    ("• Welford's online algorithm for", LIGHT, False),
    ("  running mean/std across batches", LIGHT, False),
    ("• Stats stored in SSM Parameter Store", MUTED, False),
], 10)

# Training columns callout
rect(0.6, 3.9, 7.3, 0.5, DARK, PURPLE, 1)
txt(0.9, 3.95, 6.5, 0.2, "Training columns: [pressure, temperature, delta_pressure, delta_temp]", 10, PURPLE, True)
txt(0.9, 4.15, 6.5, 0.2, "Inference columns: [vehicle_id, tire_id, first_ts, last_ts] + training columns", 10, MUTED)

# =====================================================================
# TOP RIGHT: MODEL CONFIGURATION
# =====================================================================
txt(8.3, 1.2, 7, 0.25, "MODEL CONFIGURATION", 13, ORANGE, True)

rect(8.3, 1.5, 7.1, 2.9, CARD, SKY, 1.5)

mtxt(8.6, 1.55, 3.2, 2.8, [
    ("Algorithm:", WHITE, True),
    ("SageMaker Random Cut Forest", SKY, True),
    ("", LIGHT, False),
    ("• Unsupervised anomaly detection", LIGHT, False),
    ("• No labeled training data required", LIGHT, False),
    ("• Outputs anomaly score (0.0 – 1.0)", LIGHT, False),
    ("• Learns normal patterns from fleet", LIGHT, False),
    ("", LIGHT, False),
    ("Hyperparameters:", WHITE, True),
    ("• feature_dim (number of features)", LIGHT, False),
    ("• num_samples_per_tree: 256", LIGHT, False),
    ("• num_trees: 100", LIGHT, False),
], 10)

mtxt(12.0, 1.55, 3.2, 2.8, [
    ("Training Infrastructure:", WHITE, True),
    ("• 4× ml.m5.12xlarge instances", SKY, False),
    ("• 400 GiB volume per instance", LIGHT, False),
    ("• Max runtime: 3600s (1 hour)", LIGHT, False),
    ("• RCF container from ECR", MUTED, False),
    ("", LIGHT, False),
    ("Training Pipeline (Step Function):", WHITE, True),
    ("1. ML ETL cleaner Lambda", LIGHT, False),
    ("2. SageMaker Training Job", LIGHT, False),
    ("3. Create SageMaker Model", LIGHT, False),
    ("4. Create/Update Endpoint Config", LIGHT, False),
    ("5. Update Endpoint (or Create new)", LIGHT, False),
], 10)

# =====================================================================
# BOTTOM LEFT: BATCH INFERENCE
# =====================================================================
txt(0.6, 4.65, 7, 0.25, "BATCH INFERENCE PIPELINE", 13, ORANGE, True)

rect(0.6, 4.95, 7.3, 1.7, CARD, GREEN, 1.5)

mtxt(0.9, 5.0, 3.3, 1.6, [
    ("Daily Step Function:", WHITE, True),
    ("1. Lambda reads model name from SSM", LIGHT, False),
    ("2. Starts Batch Transform Job:", LIGHT, False),
    ("   • Instance: ml.m5.12xlarge", SKY, False),
    ("   • Input: inference data bucket", LIGHT, False),
    ("   • InputFilter: features only ($[4:7])", MUTED, False),
    ("   • OutputFilter: metadata + score", MUTED, False),
], 10)

mtxt(4.2, 5.0, 3.5, 1.6, [
    ("3. Status polling (30s intervals)", LIGHT, False),
    ("4. Raw predictions → S3 (.csv.out)", LIGHT, False),
    ("", LIGHT, False),
    ("Anomaly Thresholding:", WHITE, True),
    ("• 99th percentile of batch scores", GREEN, True),
    ("• Threshold stored in SSM", LIGHT, False),
    ("• Updated each batch run", MUTED, False),
], 10)

# =====================================================================
# BOTTOM RIGHT: REAL-TIME API
# =====================================================================
txt(8.3, 4.65, 7, 0.25, "REAL-TIME PREDICTION API", 13, ORANGE, True)

rect(8.3, 4.95, 7.1, 1.7, CARD, GREEN, 1.5)

mtxt(8.6, 5.0, 3.2, 1.6, [
    ("POST /predict", GREEN, True),
    ("API Gateway → Lambda → SageMaker", LIGHT, False),
    ("", LIGHT, False),
    ("Request body:", WHITE, True),
    ("{ vehicle_id, tire_id,", LIGHT, False),
    ("  pressure, temperature,", LIGHT, False),
    ("  delta_pressure, delta_temp }", LIGHT, False),
], 10)

mtxt(12.0, 5.0, 3.2, 1.6, [
    ("Response:", GREEN, True),
    ("{ anomaly_score: 0.754,", LIGHT, False),
    ("  prediction: true }", LIGHT, False),
    ("", LIGHT, False),
    ("Flow:", WHITE, True),
    ("Normalize inputs (z-score from SSM)", LIGHT, False),
    ("→ Invoke SageMaker endpoint (CSV)", LIGHT, False),
    ("→ Score > threshold = anomaly", MUTED, False),
], 10)

# Rate limits
rect(8.3, 6.75, 7.1, 0.45, DARK, GREEN, 1)
txt(8.6, 6.8, 6.5, 0.2, "Rate: 100 req/s · Burst: 200 · Daily quota: 10K · API Key auth · CORS enabled", 10, GREEN)
txt(8.6, 7.0, 6.5, 0.2, "503 if threshold not yet established (run batch inference first)", 9, MUTED)

# =====================================================================
# BOTTOM: DUAL-PATH COMPARISON
# =====================================================================
txt(0.6, 6.95, 7, 0.25, "WHY DUAL-PATH?", 13, ORANGE, True)

rect(0.6, 7.25, 3.5, 0.85, CARD, PURPLE, 1)
mtxt(0.9, 7.3, 3, 0.75, [
    ("ML Path (Random Cut Forest)", PURPLE, True),
    ("Catches novel, unseen failure patterns", LIGHT, False),
    ("Fleet-wide learning, no domain rules", MUTED, False),
], 10)

rect(4.4, 7.25, 3.5, 0.85, CARD, YELLOW, 1)
mtxt(4.7, 7.3, 3, 0.75, [
    ("Filter Path (Physics-based)", YELLOW, True),
    ("Deterministic, explainable leak rates", LIGHT, False),
    ("Severity + time-to-threshold estimates", MUTED, False),
], 10)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide20b_pred_maint_model.pptx")
prs.save(out)
print(f"Saved: {out}")
