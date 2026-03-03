from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(s, l, t, w, h, lines):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, sz, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(2)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === TITLE ===
txt(slide, 0.6, 0.25, 12, 0.6, "COMPONENTS OF A CONNECTED VEHICLE PLATFORM", 36, ORANGE, True)
txt(slide, 0.6, 0.8, 14, 0.4, "The building blocks that power modern connected vehicle experiences", 20, LIGHT)

# === 6 COMPONENT CARDS — 3x2 grid ===
components = [
    {
        "icon": "📶",       # cellular/signal icon
        "name": "Cellular Connectivity",
        "sub": "AWS Connectivity Services (ACS)",
        "desc": "Managed cellular connectivity that gets vehicle signals to the cloud — no carrier contracts to manage",
        "color": SKY,
    },
    {
        "icon": "🔗",       # link/connection icon
        "name": "Device Connectivity",
        "sub": "AWS IoT Core",
        "desc": "Bi-directional vehicle-to-cloud communication — connect millions of devices with managed MQTT",
        "color": GREEN,
    },
    {
        "icon": "⚡",       # streaming/lightning icon
        "name": "Event Streaming",
        "sub": "Amazon MSK",
        "desc": "High-throughput event ingestion and buffering — replay, fan-out, and durable message delivery",
        "color": YELLOW,
    },
    {
        "icon": "⚙️",       # processing/gear icon
        "name": "Stream Processing",
        "sub": "Amazon Managed Flink",
        "desc": "Real-time analytics and anomaly detection on vehicle telemetry — sub-second latency at any scale",
        "color": ORANGE,
    },
    {
        "icon": "📊",       # analytics/chart icon
        "name": "Data & Analytics",
        "sub": "S3 + Athena",
        "desc": "Serverless data lake and SQL analytics — query petabytes of vehicle data without infrastructure",
        "color": PURPLE,
    },
    {
        "icon": "📞",       # engagement/phone icon
        "name": "Customer Engagement",
        "sub": "Amazon Connect",
        "desc": "AI-powered contact center and proactive outreach — triggered by vehicle events and predictions",
        "color": PINK,
    },
]

cw = 4.7
ch = 3.0
cgap_x = 0.35
cgap_y = 0.3
start_x = 0.6
start_y = 1.5

for i, comp in enumerate(components):
    col = i % 3
    row = i // 3
    x = start_x + col * (cw + cgap_x)
    y = start_y + row * (ch + cgap_y)
    color = comp["color"]

    # Card background
    rect(slide, x, y, cw, ch, CARD, color, 1.5)

    # Icon circle
    circ(slide, x + 0.3, y + 0.3, 0.8, DARK, color, 2)
    txt(slide, x + 0.3, y + 0.4, 0.8, 0.6, comp["icon"], 28, WHITE, False, PP_ALIGN.CENTER)

    # Component name
    txt(slide, x + 1.3, y + 0.3, 3.2, 0.4, comp["name"], 20, WHITE, True)

    # AWS service subheading
    txt(slide, x + 1.3, y + 0.7, 3.2, 0.3, comp["sub"], 13, color)

    # Description
    txt(slide, x + 0.3, y + 1.3, cw - 0.6, 1.2, comp["desc"], 13, LIGHT)

    # Number badge top-right
    circ(slide, x + cw - 0.55, y + 0.15, 0.4, color, None)
    txt(slide, x + cw - 0.55, y + 0.18, 0.4, 0.35, str(i+1), 14, DARK, True, PP_ALIGN.CENTER)

# === OBSERVABILITY BAR — spans bottom ===
obs_y = start_y + 2 * (ch + cgap_y) + 0.15
rect(slide, 0.6, obs_y, 15.0, 0.55, CARD, RED, 1.5)
txt(slide, 0.9, obs_y + 0.05, 0.5, 0.45, "🔍", 22, WHITE, False, PP_ALIGN.CENTER)
txt(slide, 1.5, obs_y + 0.05, 3, 0.25, "Observability", 18, WHITE, True)
txt(slide, 1.5, obs_y + 0.3, 3, 0.25, "Amazon CloudWatch", 12, RED)
txt(slide, 5.0, obs_y + 0.12, 10, 0.3, "Unified metrics, logs, alarms, and dashboards across all components — single pane of glass", 13, LIGHT)

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide6_components.pptx")
prs.save(out)
print(f"Saved: {out}")
