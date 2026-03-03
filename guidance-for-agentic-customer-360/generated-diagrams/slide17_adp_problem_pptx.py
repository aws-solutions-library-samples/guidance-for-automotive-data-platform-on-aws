from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === TITLE ===
txt(slide, 0.6, 0.2, 14, 0.6, "THE PROBLEM: AUTOMOTIVE DATA LIVES IN SILOS", 36, ORANGE, True)
txt(slide, 0.6, 0.75, 14, 0.4,
    "Each system owns a piece of the customer — none of them see the full picture", 20, LIGHT)

# === 5 SILO BOXES with red X between them ===
silos = [
    ("CRM\nPlatform", "Salesforce\nDynamics", SKY),
    ("Dealer Mgmt\nSystem", "CDK · Reynolds\n& Reynolds", GREEN),
    ("Vehicle\nTelemetry", "IoT · TCU\nCAN bus", YELLOW),
    ("Service\nRecords", "DMS · Warranty\nRecall systems", PURPLE),
    ("Contact\nCenter", "IVR · Chat\nEmail · SMS", PINK),
]

sw = 2.5
sgap = 0.35
sx = 0.6
sy = 1.5

for i, (name, systems, color) in enumerate(silos):
    x = sx + i * (sw + sgap)
    rect(slide, x, sy, sw, 1.5, CARD, color, 2)
    txt(slide, x + 0.1, sy + 0.15, sw - 0.2, 0.6, name, 14, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, x + 0.1, sy + 0.8, sw - 0.2, 0.55, systems, 11, MUTED, False, PP_ALIGN.CENTER)

    if i < len(silos) - 1:
        cx = x + sw + sgap / 2
        txt(slide, cx - 0.15, sy + 0.4, 0.3, 0.5, "✕", 22, RED, True, PP_ALIGN.CENTER)

# === QUESTIONS OEMs CAN'T ANSWER ===
txt(slide, 0.6, 3.3, 12, 0.4, "Questions OEMs can't answer today:", 18, WHITE, True)

questions = [
    ("\"Which high-value customers are at risk of churning?\"",
     "Requires CRM + telemetry + service + contact center data"),
    ("\"Can we predict and prevent failures before they strand a customer?\"",
     "Requires telemetry + service history + parts inventory + ML models"),
    ("\"What's the root cause of declining customer sentiment?\"",
     "Requires NPS + case data + vehicle health + interaction history"),
    ("\"How do we comply with EU Data Act while sharing vehicle data?\"",
     "Requires governance + PII detection + regional data residency"),
]

qy = 3.85
for question, context in questions:
    rect(slide, 0.6, qy, 14.8, 0.65, CARD)
    txt(slide, 0.9, qy + 0.05, 9, 0.3, question, 14, ORANGE, True)
    txt(slide, 0.9, qy + 0.35, 13.5, 0.25, context, 11, MUTED)
    qy += 0.75

# === PUNCHLINE ===
rect(slide, 0.6, 6.95, 14.8, 0.6, CARD, ORANGE, 2)
txt(slide, 0.8, 7.0, 14.4, 0.5,
    "You don't have a data problem — you have a data silo problem. The Automotive Data Platform solves it.",
    16, WHITE, True, PP_ALIGN.CENTER)

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide17_adp_problem.pptx")
prs.save(out)
print(f"Saved: {out}")
