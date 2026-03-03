from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)
BG = RGBColor(13, 17, 33)
PLACEHOLDER_BG = RGBColor(18, 24, 45)

def make_slide():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return prs, slide

def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(slide, l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
        if i > 0: p.space_before = Pt(2)

def rect(slide, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def placeholder(slide, l, t, w, h, label="ARCHITECTURE DIAGRAM PLACEHOLDER"):
    """Dashed-border placeholder area for a diagram."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PLACEHOLDER_BG
    sh.line.color.rgb = ORANGE; sh.line.width = Pt(1.5)
    sh.line.dash_style = 4  # dash
    txt(slide, l + 0.2, t + h/2 - 0.15, w - 0.4, 0.3, label, 12, ORANGE, False, PP_ALIGN.CENTER)

def footer(slide):
    txt(slide, 0.6, 8.6, 12, 0.3,
        "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

# =====================================================================
# VARIANT A: Large diagram top + minimal callouts bottom
# =====================================================================
prs_a, sl_a = make_slide()

txt(sl_a, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — PIPELINE ARCHITECTURE", 34, ORANGE, True)
txt(sl_a, 0.6, 0.65, 14, 0.3,
    "Dual-path tire failure prediction: ML anomaly detection + physics-based slow leak filtering", 16, LIGHT)

# Large diagram placeholder — top 60% of slide
placeholder(sl_a, 0.6, 1.15, 14.8, 4.0,
    "← INSERT ARCHITECTURE DIAGRAM: Redshift → Root ETL → [ML Path | Filter Path] → Alerts →")

# Below diagram: clean text callouts, no boxes — just colored headers + text
txt(sl_a, 0.6, 5.4, 3, 0.25, "DATA SOURCE", 11, SKY, True)
mtxt(sl_a, 0.6, 5.7, 3.5, 0.9, [
    ("Redshift Datashare (MMT telemetry)", WHITE, False),
    ("Hourly SQL query via Lambda → S3", LIGHT, False),
    ("Fields: pressure, temperature,", LIGHT, False),
    ("vehicle_id, tire_id, timestamp", MUTED, False),
], 10)

txt(sl_a, 4.3, 5.4, 3, 0.25, "ROOT ETL", 11, GREEN, True)
mtxt(sl_a, 4.3, 5.7, 3.5, 0.9, [
    ("Glue Spark job (hourly, 30min offset)", WHITE, False),
    ("Clean · convert units · weave tables", LIGHT, False),
    ("Output: S3 partitioned by date/hour", LIGHT, False),
    ("Feeds both ML and Filtering paths", MUTED, False),
], 10)

txt(sl_a, 8.0, 5.4, 3, 0.25, "ML PATH", 11, PURPLE, True)
mtxt(sl_a, 8.0, 5.7, 3.5, 0.9, [
    ("Daily ML ETL → resample to 1-day", WHITE, False),
    ("Engineer: delta_pressure, delta_temp", LIGHT, False),
    ("Normalize: z-score (Welford's → SSM)", LIGHT, False),
    ("SageMaker RCF → anomaly scores", MUTED, False),
], 10)

txt(sl_a, 11.7, 5.4, 3, 0.25, "FILTERING PATH", 11, YELLOW, True)
mtxt(sl_a, 11.7, 5.7, 3.5, 0.9, [
    ("Nightly Step Function → Glue job", WHITE, False),
    ("Quantile filter (10th pctl rolling)", LIGHT, False),
    ("Leak rate = ΔPSI / Δdays", LIGHT, False),
    ("Severity: H>3 · M>2 · L>1 PSI/day", MUTED, False),
], 10)

# Thin accent line separator
for x in range(60, 1540):
    pass  # skip pixel drawing in pptx — use a shape instead
sh = sl_a.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.85), Inches(14.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(35, 45, 72); sh.line.fill.background()

# Alerts row — single line, no box
txt(sl_a, 0.6, 7.0, 2, 0.25, "ALERTS SYSTEM", 11, ORANGE, True)
mtxt(sl_a, 2.8, 7.0, 12, 0.5, [
    ("Both paths → S3 upload triggers alerts-transformer Lambda → DynamoDB (PENDING) → SQS → processor (PROCESSED) → DLQ → cleaner (FAILED)", LIGHT, False),
    ("ML threshold: 99th percentile of batch scores, stored in SSM · Filtering: severity classification + time-to-80-PSI estimate", MUTED, False),
], 9)

# Metrics as inline text, not boxes
sh = sl_a.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.75), Inches(14.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
txt(sl_a, 0.6, 7.85, 14.8, 0.3,
    "7-14 day advance warning  ·  4 features (pressure, temp, Δp, Δt)  ·  Dual-path validation  ·  99th percentile threshold  ·  3 severity levels  ·  DynamoDB alert lifecycle",
    10, ORANGE)

footer(sl_a)
prs_a.save(os.path.join(OUTDIR, "slide20a_pipeline_varA.pptx"))
print("Saved: slide20a_pipeline_varA.pptx")

# =====================================================================
# VARIANT B: Left diagram placeholder + right clean text columns
# =====================================================================
prs_b, sl_b = make_slide()

txt(sl_b, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — PIPELINE ARCHITECTURE", 34, ORANGE, True)
txt(sl_b, 0.6, 0.65, 14, 0.3,
    "Dual-path tire failure prediction: ML anomaly detection + physics-based slow leak filtering", 16, LIGHT)

# Left: diagram placeholder (60% width)
placeholder(sl_b, 0.6, 1.15, 9.0, 5.8,
    "← INSERT ARCHITECTURE DIAGRAM →")

# Right: clean text — no boxes, just headers + body
x = 10.0

txt(sl_b, x, 1.15, 5.5, 0.25, "① DATA SOURCE", 12, SKY, True)
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.42), Inches(1.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = SKY; sh.line.fill.background()
mtxt(sl_b, x, 1.5, 5.3, 0.5, [
    ("Redshift Datashare — MMT tire telemetry", WHITE, False),
    ("Hourly Lambda → SQL query → raw CSV to S3", LIGHT, False),
], 10)

txt(sl_b, x, 2.15, 5.5, 0.25, "② ROOT ETL", 12, GREEN, True)
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.42), Inches(1.5), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = GREEN; sh.line.fill.background()
mtxt(sl_b, x, 2.5, 5.3, 0.7, [
    ("Glue Spark (hourly, 30min offset from query)", WHITE, False),
    ("Clean erroneous values · unit conversions", LIGHT, False),
    ("Weave multi-table by vehicle_id + timestamp", LIGHT, False),
    ("→ S3 partitioned by date/hour", MUTED, False),
], 10)

txt(sl_b, x, 3.4, 5.5, 0.25, "③ ML PATH (daily)", 12, PURPLE, True)
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.67), Inches(2.0), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = PURPLE; sh.line.fill.background()
mtxt(sl_b, x, 3.75, 5.3, 1.0, [
    ("ML ETL: resample 1-day · engineer delta_pressure,", WHITE, False),
    ("delta_temp · normalize z-score (Welford's → SSM)", LIGHT, False),
    ("Training: SageMaker RCF · 4× ml.m5.12xlarge", LIGHT, False),
    ("Inference: daily batch transform · anomaly scores", LIGHT, False),
    ("Threshold: 99th percentile, stored in SSM", MUTED, False),
], 10)

txt(sl_b, x, 5.0, 5.5, 0.25, "③ FILTERING PATH (nightly)", 12, YELLOW, True)
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.27), Inches(3.0), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = YELLOW; sh.line.fill.background()
mtxt(sl_b, x, 5.35, 5.3, 0.9, [
    ("Step Function → Lambda → Glue job per tire", WHITE, False),
    ("Quantile filter (10th pctl rolling window)", LIGHT, False),
    ("Leak rate = ΔPSI / Δdays", LIGHT, False),
    ("Severity: HIGH >3 · MED >2 · LOW >1 PSI/day", LIGHT, False),
    ("Time-to-80-PSI threshold estimate", MUTED, False),
], 10)

txt(sl_b, x, 6.5, 5.5, 0.25, "④ ALERTS", 12, ORANGE, True)
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(6.77), Inches(1.2), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
mtxt(sl_b, x, 6.85, 5.3, 0.5, [
    ("S3 trigger → DynamoDB (PENDING) → SQS →", WHITE, False),
    ("processor (PROCESSED) · DLQ → cleaner (FAILED)", LIGHT, False),
], 10)

# Bottom metrics line
sh = sl_b.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.75), Inches(14.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
txt(sl_b, 0.6, 7.85, 14.8, 0.3,
    "7-14 day advance warning  ·  4 features  ·  Dual-path (ML + physics)  ·  99th pctl threshold  ·  3 severity levels  ·  DynamoDB lifecycle",
    10, ORANGE)

footer(sl_b)
prs_b.save(os.path.join(OUTDIR, "slide20a_pipeline_varB.pptx"))
print("Saved: slide20a_pipeline_varB.pptx")

# =====================================================================
# VARIANT C: Top diagram + bottom two-column (ML vs Filter) clean text
# =====================================================================
prs_c, sl_c = make_slide()

txt(sl_c, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — PIPELINE ARCHITECTURE", 34, ORANGE, True)
txt(sl_c, 0.6, 0.65, 14, 0.3,
    "Dual-path tire failure prediction: ML anomaly detection + physics-based slow leak filtering", 16, LIGHT)

# Top: diagram placeholder
placeholder(sl_c, 0.6, 1.15, 14.8, 3.2,
    "← INSERT ARCHITECTURE DIAGRAM: Redshift → Root ETL → [ML | Filter] → Alerts →")

# Shared pipeline intro — no box, just text
txt(sl_c, 0.6, 4.55, 14, 0.25, "SHARED PIPELINE", 11, GREEN, True)
sh = sl_c.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.82), Inches(14.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(35, 45, 72); sh.line.fill.background()
txt(sl_c, 0.6, 4.9, 14.8, 0.3,
    "Redshift Datashare (hourly query) → Root ETL Glue Spark (clean, convert, weave multi-table by vehicle_id + timestamp) → S3 partitioned by date/hour",
    10, LIGHT)

# Two columns below: ML vs Filtering — clean text, one accent box each
# Left: ML
txt(sl_c, 0.6, 5.4, 7, 0.25, "ML ANOMALY DETECTION", 13, PURPLE, True)
sh = sl_c.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.67), Inches(0.06), Inches(2.0))
sh.fill.solid(); sh.fill.fore_color.rgb = PURPLE; sh.line.fill.background()

mtxt(sl_c, 0.9, 5.7, 6.8, 2.0, [
    ("ML ETL (daily Step Function → Glue)", WHITE, True),
    ("Resample to 1-day intervals · Engineer delta_pressure, delta_temp", LIGHT, False),
    ("Normalize via z-score using Welford's online algorithm (stats in SSM)", LIGHT, False),
    ("", LIGHT, False),
    ("Training (Step Function → SageMaker)", WHITE, True),
    ("Random Cut Forest · 4× ml.m5.12xlarge · 400 GiB · Unsupervised", LIGHT, False),
    ("Hyperparams: feature_dim, num_samples_per_tree (256), num_trees (100)", MUTED, False),
    ("", LIGHT, False),
    ("Inference (daily Step Function)", WHITE, True),
    ("Batch Transform (ml.m5.12xlarge) → anomaly score per tire per day", LIGHT, False),
    ("Threshold: 99th percentile of batch scores → SSM", MUTED, False),
], 10)

# Right: Filtering
txt(sl_c, 8.3, 5.4, 7, 0.25, "PHYSICS-BASED LEAK DETECTION", 13, YELLOW, True)
sh = sl_c.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(5.67), Inches(0.06), Inches(2.0))
sh.fill.solid(); sh.fill.fore_color.rgb = YELLOW; sh.line.fill.background()

mtxt(sl_c, 8.6, 5.7, 6.8, 2.0, [
    ("Detection (nightly Step Function → Glue)", WHITE, True),
    ("Quantile filter: 10th percentile rolling window", LIGHT, False),
    ("Daily aggregation to smooth sensor noise", LIGHT, False),
    ("Detect pressure drops exceeding threshold over time window", LIGHT, False),
    ("", LIGHT, False),
    ("Leak Rate Calculation", WHITE, True),
    ("(first_psi − last_psi) / Δdays = PSI/day", LIGHT, False),
    ("Severity: HIGH >3 · MEDIUM >2 · LOW >1 PSI/day", LIGHT, False),
    ("Time-to-80-PSI threshold estimate for maintenance planning", MUTED, False),
    ("", LIGHT, False),
    ("CloudWatch: num_slow_leak_detected, num_aaids, processing_time_s", MUTED, False),
], 10)

# Alerts — single line
sh = sl_c.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.85), Inches(14.8), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
txt(sl_c, 0.6, 7.55, 2, 0.25, "ALERTS", 11, ORANGE, True)
txt(sl_c, 2.2, 7.55, 13, 0.25,
    "Both paths → S3 trigger → DynamoDB (PENDING) → SQS → processor (PROCESSED) → DLQ → cleaner (FAILED) · 99th pctl threshold",
    10, LIGHT)

txt(sl_c, 0.6, 7.9, 14.8, 0.3,
    "7-14 day advance warning  ·  4 features  ·  Dual-path validation  ·  Real-time API + batch  ·  3 severity levels",
    10, ORANGE)

footer(sl_c)
prs_c.save(os.path.join(OUTDIR, "slide20a_pipeline_varC.pptx"))
print("Saved: slide20a_pipeline_varC.pptx")
