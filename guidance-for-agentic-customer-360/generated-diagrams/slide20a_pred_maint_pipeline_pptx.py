from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — PIPELINE ARCHITECTURE", 34, ORANGE, True)
txt(0.6, 0.65, 14, 0.35,
    "From Redshift telemetry to dual-path prediction: ML anomaly detection + physics-based leak filtering", 17, LIGHT)

# =====================================================================
# FULL-WIDTH PIPELINE FLOW
# =====================================================================

# --- ① Data Source ---
rect(0.6, 1.25, 14.8, 0.7, CARD, SKY, 1.5)
txt(0.9, 1.3, 3, 0.25, "① Data Source", 14, SKY, True)
txt(4.0, 1.3, 10, 0.25, "Redshift Datashare — MMT tire telemetry (pressure, temperature, vehicle_id, tire_id, timestamp)", 11, LIGHT)
txt(0.9, 1.6, 14, 0.25, "Hourly CloudWatch cron → Redshift Query Lambda → SQL query on datashare tables → raw CSV uploaded to S3 raw data bucket", 10, MUTED)

arrow(8.0, 1.95, 8.0, 2.2, SKY, 1.5)

# --- ② Root ETL ---
rect(0.6, 2.2, 14.8, 0.85, CARD, GREEN, 1.5)
txt(0.9, 2.25, 3, 0.25, "② Root ETL", 14, GREEN, True)
txt(4.0, 2.25, 4, 0.25, "Amazon Glue (Spark) · Hourly · 30 min offset from query", 11, LIGHT)
mtxt(0.9, 2.55, 14, 0.45, [
    ("Clean erroneous values · Parse timestamps (multi-format) · Unit conversions · Drop null vehicle_id/tire_id/timestamp", LIGHT, False),
    ("Weave data from multiple Redshift tables into single source by vehicle_id + timestamp → Output: S3 ETL bucket, partitioned by date/hour", MUTED, False),
], 10)

# Split arrows
arrow(5.0, 3.05, 5.0, 3.35, GREEN, 1.5)
arrow(11.0, 3.05, 11.0, 3.35, GREEN, 1.5)
txt(3.5, 3.1, 3, 0.25, "ML PATH", 10, PURPLE, True, PP_ALIGN.CENTER)
txt(9.5, 3.1, 3, 0.25, "FILTERING PATH", 10, YELLOW, True, PP_ALIGN.CENTER)

# --- ③ ML ETL (left) ---
rect(0.6, 3.35, 7.0, 1.5, CARD, PURPLE, 1.5)
txt(0.9, 3.4, 6, 0.25, "③ ML ETL Pipeline (Daily Step Function)", 13, PURPLE, True)
mtxt(0.9, 3.7, 3.2, 1.1, [
    ("Lambda triggers Glue Job:", WHITE, False),
    ("1. Add metadata (first/last timestamp)", LIGHT, False),
    ("2. Resample to 1-day intervals", LIGHT, False),
    ("   Group by vehicle_id + tire_id", MUTED, False),
    ("   Apply mean/median/mode", MUTED, False),
], 10)
mtxt(4.2, 3.7, 3.2, 1.1, [
    ("3. Engineer features:", WHITE, False),
    ("   • delta_pressure (lag difference)", LIGHT, False),
    ("   • delta_temp (lag difference)", LIGHT, False),
    ("4. Normalize (z-score):", WHITE, False),
    ("   Welford's online algorithm → SSM", MUTED, False),
], 10)

# --- ③ Filtering (right) ---
rect(8.4, 3.35, 7.0, 1.5, CARD, YELLOW, 1.5)
txt(8.7, 3.4, 6, 0.25, "③ Slow Leak Detection (Nightly Step Function)", 13, YELLOW, True)
mtxt(8.7, 3.7, 3.2, 1.1, [
    ("Lambda → Glue Job per tire:", WHITE, False),
    ("1. Quantile filter (10th pctl rolling)", LIGHT, False),
    ("   Weights toward lower values", MUTED, False),
    ("2. Daily aggregation (smooth noise)", LIGHT, False),
    ("3. Detect pressure drops > threshold", LIGHT, False),
], 10)
mtxt(12.0, 3.7, 3.2, 1.1, [
    ("4. Calculate leak rate:", WHITE, False),
    ("   (first_psi − last_psi) / Δdays", LIGHT, False),
    ("5. Classify severity:", WHITE, False),
    ("   HIGH > 3 · MED > 2 · LOW > 1 PSI/day", GREEN, False),
    ("6. Estimate time to 80 PSI threshold", MUTED, False),
], 10)

# Arrows down
arrow(4.1, 4.85, 4.1, 5.15, PURPLE, 1.5)
arrow(11.9, 4.85, 11.9, 5.15, YELLOW, 1.5)

# --- ④ ML Inference (left) ---
rect(0.6, 5.15, 7.0, 1.15, CARD, PURPLE, 1.5)
txt(0.9, 5.2, 6, 0.25, "④ ML Inference (Daily Step Function)", 13, PURPLE, True)
mtxt(0.9, 5.5, 6.5, 0.75, [
    ("Lambda reads model name from SSM → starts SageMaker Batch Transform (ml.m5.12xlarge)", LIGHT, False),
    ("Input: normalized features from ML ETL bucket · Output: anomaly score per vehicle_id + tire_id", LIGHT, False),
    ("Status polling Lambda (30s intervals) → on completion, raw predictions (.csv.out) written to S3", MUTED, False),
], 10)

# --- ④ Filtering Output (right) ---
rect(8.4, 5.15, 7.0, 1.15, CARD, YELLOW, 1.5)
txt(8.7, 5.2, 6, 0.25, "④ Filtering Results", 13, YELLOW, True)
mtxt(8.7, 5.5, 6.5, 0.75, [
    ("Per tire: leak_rate (PSI/day), severity (HIGH/MED/LOW), time_to_80_psi (days)", LIGHT, False),
    ("Results written to S3 as CSV · Optional DynamoDB table (id + aaid sort key)", LIGHT, False),
    ("CloudWatch metrics: num_slow_leak_detected, num_aaids, num_tires, processing_time_s", MUTED, False),
], 10)

# Merge arrows into Alerts
arrow(4.1, 6.3, 8.0, 6.6, ORANGE, 1.5)
arrow(11.9, 6.3, 8.0, 6.6, ORANGE, 1.5)

# --- ⑤ Alerts System ---
rect(0.6, 6.55, 14.8, 1.1, CARD, ORANGE, 2)
txt(0.9, 6.6, 6, 0.25, "⑤ Unified Alerts System", 14, ORANGE, True)
txt(8.0, 6.6, 7, 0.25, "Both ML and Filtering results converge here", 11, MUTED)
mtxt(0.9, 6.9, 14, 0.7, [
    ("S3 upload (ml/ or filtering/ prefix) → S3 Event Notification → alerts-transformer Lambda", LIGHT, False),
    ("→ Writes PENDING alerts to DynamoDB (partition: alertId = source_vehicleId_tireId, sort: timestamp) → sends to SQS alerts-processing queue", LIGHT, False),
    ("→ alerts-processor Lambda: reads DynamoDB, updates status to PROCESSED, writes downstream CSV · Failed → DLQ → alerts-cleaner sets FAILED", MUTED, False),
], 10)

# === BOTTOM: Key Metrics Bar ===
rect(0.5, 7.9, 15.0, 0.45, CARD, ORANGE, 1)
metrics = [
    ("Hourly", "data ingestion", 0.8),
    ("Daily", "ML ETL + inference", 2.8),
    ("Nightly", "leak detection", 4.8),
    ("Dual-path", "ML + physics", 6.8),
    ("3 severity", "HIGH / MED / LOW", 9.0),
    ("DynamoDB", "alert lifecycle", 11.2),
    ("CloudWatch", "full observability", 13.2),
]
for val, label, x in metrics:
    txt(x, 7.93, 2, 0.2, val, 11, ORANGE, True)
    txt(x, 8.12, 2, 0.2, label, 8, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide20a_pred_maint_pipeline.pptx")
prs.save(out)
print(f"Saved: {out}")
