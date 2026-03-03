from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import math, os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def icon_shape(s, l, t, sz, shape_type):
    """Placeholder geometric shape — swap for real icons in PowerPoint/Canva."""
    sh = s.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.background()
    sh.line.color.rgb = WHITE
    sh.line.width = Pt(1.5)

# === TITLE ===
txt(slide, 0.6, 0.25, 14, 0.6, "ELEMENTS OF A CONNECTED VEHICLE PLATFORM", 36, ORANGE, True)
txt(slide, 0.6, 0.8, 14, 0.4, "Six managed AWS services powering modern connected vehicle experiences", 20, LIGHT)

# === CIRCLE LAYOUT — 6 nodes around a ring ===
cx, cy = 5.2, 4.8
radius = 2.9
node_sz = 1.7
icon_sz = 0.4

elements = [
    ("Cellular\nConnectivity", "ACS", SKY, MSO_SHAPE.DIAMOND),
    ("Device\nConnectivity", "IoT Core", GREEN, MSO_SHAPE.OVAL),
    ("Event\nStreaming", "MSK", YELLOW, MSO_SHAPE.CHEVRON),
    ("Stream\nProcessing", "Flink", ORANGE, MSO_SHAPE.HEXAGON),
    ("Data &\nAnalytics", "S3 + Athena", PURPLE, MSO_SHAPE.PENTAGON),
    ("Customer\nEngagement", "Connect", PINK, MSO_SHAPE.HEART),
]

# Calculate positions (top-center first, clockwise)
positions = []
for i in range(6):
    angle = math.radians(-90 + i * 60)
    x = cx + radius * math.cos(angle) - node_sz / 2
    y = cy + radius * math.sin(angle) - node_sz / 2
    positions.append((x, y))

# Dashed connector lines between adjacent nodes
for i in range(6):
    ni = (i + 1) % 6
    x1 = positions[i][0] + node_sz / 2
    y1 = positions[i][1] + node_sz / 2
    x2 = positions[ni][0] + node_sz / 2
    y2 = positions[ni][1] + node_sz / 2
    dx, dy = x2 - x1, y2 - y1
    dist = math.sqrt(dx*dx + dy*dy)
    off = node_sz * 0.52
    sx, sy = x1 + dx/dist * off, y1 + dy/dist * off
    ex, ey = x2 - dx/dist * off, y2 - dy/dist * off
    conn = slide.shapes.add_connector(1, Inches(sx), Inches(sy), Inches(ex), Inches(ey))
    conn.line.color.rgb = RGBColor(40, 52, 75)
    conn.line.width = Pt(1.5)
    conn.line.dash_style = 2

# Draw element nodes
for i, (name, service, color, shape) in enumerate(elements):
    x, y = positions[i]

    # Outer circle with accent border
    circ(slide, x, y, node_sz, CARD, color, 2)

    # Placeholder icon centered in upper area
    ix = x + (node_sz - icon_sz) / 2
    iy = y + 0.22
    icon_shape(slide, ix, iy, icon_sz, shape)

    # Element name
    txt(slide, x + 0.05, y + 0.72, node_sz - 0.1, 0.55, name, 12, WHITE, True, PP_ALIGN.CENTER)

    # AWS service label
    txt(slide, x + 0.05, y + 1.25, node_sz - 0.1, 0.3, service, 10, color, False, PP_ALIGN.CENTER)

    # Number badge
    circ(slide, x - 0.1, y - 0.1, 0.35, color, None)
    txt(slide, x - 0.1, y - 0.07, 0.35, 0.3, str(i + 1), 13, DARK, True, PP_ALIGN.CENTER)

# === CENTER LABEL ===
circ(slide, cx - 1.0, cy - 1.0, 2.0, DARK, ORANGE, 2.5)
txt(slide, cx - 0.9, cy - 0.5, 1.8, 0.35, "Connected", 18, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.9, cy - 0.1, 1.8, 0.35, "Vehicle", 18, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.9, cy + 0.3, 1.8, 0.35, "Platform", 18, ORANGE, True, PP_ALIGN.CENTER)

# === RIGHT PANEL — element descriptions ===
rx = 10.2
rect(slide, rx, 1.4, 5.2, 6.8, CARD)
txt(slide, rx + 0.3, 1.5, 4.5, 0.4, "Platform Elements", 18, WHITE, True)

descs = [
    ("1", "Cellular Connectivity", "Managed cellular — vehicle signals to cloud, no carrier contracts", SKY),
    ("2", "Device Connectivity", "Bi-directional MQTT — connect and command millions of vehicles", GREEN),
    ("3", "Event Streaming", "High-throughput ingestion with durable replay and fan-out", YELLOW),
    ("4", "Stream Processing", "Real-time analytics and anomaly detection at any scale", ORANGE),
    ("5", "Data & Analytics", "Serverless data lake — query petabytes with standard SQL", PURPLE),
    ("6", "Customer Engagement", "AI-powered contact center triggered by vehicle events", PINK),
]

dy = 2.1
for num, title, desc, color in descs:
    circ(slide, rx + 0.25, dy + 0.05, 0.3, color, None)
    txt(slide, rx + 0.25, dy + 0.07, 0.3, 0.25, num, 11, DARK, True, PP_ALIGN.CENTER)
    txt(slide, rx + 0.7, dy, 4.3, 0.25, title, 14, WHITE, True)
    txt(slide, rx + 0.7, dy + 0.28, 4.3, 0.4, desc, 11, LIGHT)
    dy += 0.75

# Observability bar
rect(slide, rx + 0.2, dy + 0.15, 4.8, 0.45, DARK, RED, 1)
txt(slide, rx + 0.4, dy + 0.2, 4.4, 0.35, "Observability — Amazon CloudWatch across all layers", 12, RED)

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide6_elements.pptx")
prs.save(out)
print(f"Saved: {out}")
