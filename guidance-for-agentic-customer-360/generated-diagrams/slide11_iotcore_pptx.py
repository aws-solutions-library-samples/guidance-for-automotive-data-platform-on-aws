from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def bullet_list(s, l, t, w, h, items, sz=13, color=LIGHT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = False
        p.space_after = Pt(6)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === HEADER ===
circ(slide, 0.5, 0.2, 0.55, GREEN, None)
txt(slide, 0.5, 0.24, 0.55, 0.45, "2", 22, DARK, True, PP_ALIGN.CENTER)
txt(slide, 1.2, 0.2, 12, 0.6, "DEVICE CONNECTIVITY", 36, ORANGE, True)
txt(slide, 1.2, 0.75, 12, 0.4, "AWS IoT Core — Managed MQTT for Connected Vehicles", 22, GREEN, True)

# === LEFT COLUMN — What it does ===
lx = 0.6
txt(slide, lx, 1.4, 7, 0.35, "What IoT Core Does", 18, WHITE, True)

rect(slide, lx, 1.85, 7.0, 4.6, CARD, GREEN, 1.5)

bullet_list(slide, lx + 0.25, 1.95, 6.5, 4.4, [
    ("Fully managed MQTT message broker that securely connects", WHITE),
    ("millions of vehicles and routes trillions of messages to AWS", LIGHT),
    ("", LIGHT),
    ("▸  Bi-directional vehicle-to-cloud communication using", LIGHT),
    ("   MQTT 5 with full backwards compatibility to MQTT 3.1.1", LIGHT),
    ("", LIGHT),
    ("▸  mTLS authentication with X.509 certificates — no AWS-", LIGHT),
    ("   specific software required onboard (any MQTT client works)", LIGHT),
    ("", LIGHT),
    ("▸  Multi-region deployment with Route 53 geo-routing for", LIGHT),
    ("   global fleets and local data residency compliance", LIGHT),
    ("", LIGHT),
    ("▸  Shared subscriptions for load-balanced message processing,", LIGHT),
    ("   request/response patterns, message expiry, retained messages", LIGHT),
    ("", LIGHT),
    ("▸  IoT Rules Engine routes messages to 15+ AWS services —", LIGHT),
    ("   S3, Kinesis, Lambda, Timestream, and more", LIGHT),
], 13, LIGHT)

# === RIGHT COLUMN — Why IoT Core ===
rx = 7.9
txt(slide, rx, 1.4, 7, 0.35, "Why IoT Core", 18, WHITE, True)

rect(slide, rx, 1.85, 7.5, 4.6, CARD, ORANGE, 1.5)

bullet_list(slide, rx + 0.25, 1.95, 7.0, 4.4, [
    ("Proven at automotive scale", WHITE),
    ("Toyota, Honda, Mercedes-Benz (20M+ vehicles), WirelessCar", LIGHT),
    ("(11M vehicles) — all run production platforms on IoT Core", LIGHT),
    ("", LIGHT),
    ("Build to a specification, not a service", WHITE),
    ("Standard MQTT protocol means no vendor lock-in — any open-", LIGHT),
    ("source client (Eclipse Paho, Mosquitto) on any OS (QNX, Linux)", LIGHT),
    ("", LIGHT),
    ("Scales without capacity planning", WHITE),
    ("Auto-scales to millions of concurrent connections with 99.9%+", LIGHT),
    ("SLA — no brokers to provision, patch, or manage", LIGHT),
    ("", LIGHT),
    ("Purpose-built for automotive", WHITE),
    ("AWS co-developed MQTT 5 features with Toyota — TLS 1.3,", LIGHT),
    ("large payload support, static IPs, and dedicated device", LIGHT),
    ("gateways built from OEM requirements (Mercedes-Benz)", LIGHT),
], 13, LIGHT)

# === BOTTOM — Companion services ===
py = 6.7
txt(slide, 0.6, py, 4, 0.3, "Companion Services", 14, MUTED, True)

products = [
    ("IoT FleetWise", GREEN, "Intelligent vehicle data collection & campaigns"),
    ("IoT Device Mgmt", SKY, "OTA jobs, fleet indexing, device shadows"),
    ("IoT Greengrass", PURPLE, "Edge runtime for local compute & ML inference"),
]

px = 0.6
for name, color, desc in products:
    rect(slide, px, py + 0.35, 4.8, 0.7, CARD, color, 1)
    txt(slide, px + 0.2, py + 0.38, 2.5, 0.3, name, 14, WHITE, True)
    txt(slide, px + 0.2, py + 0.65, 4.4, 0.3, desc, 11, LIGHT)
    px += 5.0

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide11_iotcore.pptx")
prs.save(out)
print(f"Saved: {out}")
