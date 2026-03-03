from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.5, "DATA MESH FOUNDATION & GOVERNANCE", 34, ORANGE, True)
txt(0.6, 0.65, 14, 0.35,
    "SageMaker Unified Studio (DataZone V2) with federated catalog, blueprints, and domain-oriented ownership", 17, LIGHT)

# =====================================================================
# TOP: SageMaker Unified Studio Domain
# =====================================================================

rect(0.6, 1.2, 14.8, 1.6, CARD, ORANGE, 2)
txt(0.9, 1.25, 6, 0.25, "SageMaker Unified Studio Domain", 15, ORANGE, True)
txt(8.0, 1.25, 7, 0.25, "DataZone V2 · SSO via IAM Identity Center · AUTOMATIC user assignment", 10, MUTED)

mtxt(0.9, 1.55, 4.5, 1.2, [
    ("Domain: automotive-data-platform", WHITE, True),
    ("• DomainExecutionRole trusts:", LIGHT, False),
    ("  datazone, sagemaker, lakeformation, glue", SKY, False),
    ("• S3, Glue, Athena, Redshift, Lake Formation", LIGHT, False),
    ("  access policies built in", MUTED, False),
    ("• Portal URL for web-based access", MUTED, False),
], 9)

mtxt(5.6, 1.55, 4.5, 1.2, [
    ("Enabled Blueprints:", WHITE, True),
    ("• ML Analysis (SageMaker)", SKY, False),
    ("• Tooling (Glue, Athena)", GREEN, False),
    ("• Redshift (data warehouse)", PURPLE, False),
    ("• EMR (big data processing)", YELLOW, False),
    ("• Bedrock (generative AI)", PINK, False),
], 9)

mtxt(10.3, 1.55, 5, 1.2, [
    ("Project Profiles:", WHITE, True),
    ("• Default profile for auto use cases", LIGHT, False),
    ("• Team-based access control", LIGHT, False),
    ("• Shared notebooks and queries", LIGHT, False),
    ("• Data lineage tracking", MUTED, False),
    ("• Searchable data catalog", MUTED, False),
], 9)

# =====================================================================
# MIDDLE LEFT: Infrastructure
# =====================================================================

txt(0.6, 3.05, 7, 0.25, "SHARED INFRASTRUCTURE", 13, ORANGE, True)

rect(0.6, 3.35, 7.3, 1.6, CARD, SKY, 1.5)

mtxt(0.9, 3.4, 3.3, 1.5, [
    ("Networking:", WHITE, True),
    ("• VPC (10.0.0.0/16)", LIGHT, False),
    ("• 3 private subnets across AZs", LIGHT, False),
    ("• NAT Gateway for outbound", LIGHT, False),
    ("• VPC endpoints: S3, Glue,", LIGHT, False),
    ("  SageMaker, Athena, Redshift", MUTED, False),
], 9)

mtxt(4.2, 3.4, 3.5, 1.5, [
    ("Storage:", WHITE, True),
    ("• DataZone environment bucket", LIGHT, False),
    ("  (encrypted, versioned, logged)", MUTED, False),
    ("• Access logs bucket", LIGHT, False),
    ("• Glue database for catalog", LIGHT, False),
    ("• S3 data lake (per-solution)", MUTED, False),
], 9)

# =====================================================================
# MIDDLE RIGHT: Data Publishing
# =====================================================================

txt(8.3, 3.05, 7, 0.25, "DATA CATALOG & PUBLISHING", 13, ORANGE, True)

rect(8.3, 3.35, 7.1, 1.6, CARD, GREEN, 1.5)

mtxt(8.6, 3.4, 3.2, 1.5, [
    ("Glue → DataZone Publishing:", WHITE, True),
    ("• Discover Glue databases", LIGHT, False),
    ("• Create data source in DataZone", LIGHT, False),
    ("• Run discovery job", LIGHT, False),
    ("• Accept discovered assets", LIGHT, False),
    ("• Publish to catalog", MUTED, False),
], 9)

mtxt(12.0, 3.4, 3.2, 1.5, [
    ("Automotive Data Assets:", WHITE, True),
    ("• Customer 360 (CRM + service)", SKY, False),
    ("• Vehicle telemetry (IoT)", GREEN, False),
    ("• Tire prediction data", PURPLE, False),
    ("• Sales & inventory", YELLOW, False),
    ("• Weather & external data", MUTED, False),
], 9)

# =====================================================================
# BOTTOM LEFT: IAM & Roles
# =====================================================================

txt(0.6, 5.2, 7, 0.25, "IAM ROLES & ACCESS CONTROL", 13, ORANGE, True)

rect(0.6, 5.5, 7.3, 1.5, CARD, PURPLE, 1.5)

mtxt(0.9, 5.55, 3.3, 1.4, [
    ("DomainExecutionRole:", WHITE, True),
    ("• SageMaker: Describe, Create, Delete", LIGHT, False),
    ("• S3: Read/Write datazone-* buckets", LIGHT, False),
    ("• Glue: Full catalog access", LIGHT, False),
    ("• Athena: Query execution", LIGHT, False),
    ("• Lake Formation: Grant/Revoke", LIGHT, False),
    ("• Redshift: Describe + credentials", MUTED, False),
], 9)

mtxt(4.2, 5.55, 3.5, 1.4, [
    ("Blueprint Provisioning Roles:", WHITE, True),
    ("• ManageAccessRole (per blueprint)", LIGHT, False),
    ("• ProvisioningRole (environment)", LIGHT, False),
    ("• Scoped to specific blueprint ARNs", MUTED, False),
    ("", LIGHT, False),
    ("DomainServiceRole:", WHITE, True),
    ("• DataZone + SSO + IdentityStore", LIGHT, False),
], 9)

# =====================================================================
# BOTTOM RIGHT: Use Cases
# =====================================================================

txt(8.3, 5.2, 7, 0.25, "PERSONA-BASED WORKFLOWS", 13, ORANGE, True)

rect(8.3, 5.5, 3.4, 1.5, CARD, SKY, 1.5)
txt(8.6, 5.55, 3, 0.2, "Data Engineer", 12, SKY, True)
mtxt(8.6, 5.8, 3, 1.1, [
    ("• Glue ETL jobs for data prep", LIGHT, False),
    ("• Athena queries for validation", LIGHT, False),
    ("• Publish datasets to catalog", LIGHT, False),
    ("• Manage data quality", MUTED, False),
], 9)

rect(12.0, 5.5, 3.4, 1.5, CARD, PURPLE, 1.5)
txt(12.3, 5.55, 3, 0.2, "ML Engineer", 12, PURPLE, True)
mtxt(12.3, 5.8, 3, 1.1, [
    ("• SageMaker notebooks", LIGHT, False),
    ("• Model training & deployment", LIGHT, False),
    ("• Tire prediction project", LIGHT, False),
    ("• Experiment tracking", MUTED, False),
], 9)

# === BOTTOM: Architecture Summary Bar ===
rect(0.5, 7.25, 14.9, 0.5, CARD, ORANGE, 1)
txt(0.8, 7.28, 14, 0.2,
    "CloudFormation stacks: shared-resources.yaml → datazone-domain.yaml → blueprint-roles.yaml → enable-all-blueprints.yaml → policy-grants.yaml",
    10, ORANGE)
txt(0.8, 7.5, 14, 0.2,
    "Deployment: deploy-complete-platform.sh automates full stack · publish-data-sources.sh catalogs Glue databases into DataZone",
    9, MUTED)

# === BOTTOM: Key Metrics Bar ===
rect(0.5, 7.95, 14.9, 0.4, CARD, ORANGE, 1)
metrics = [
    ("DataZone V2", "unified catalog", 0.8),
    ("SSO", "IAM Identity Center", 3.0),
    ("5 blueprints", "ML, Glue, Redshift, EMR, Bedrock", 5.2),
    ("3 AZ VPC", "private subnets + endpoints", 8.2),
    ("Lake Formation", "fine-grained access", 11.2),
    ("CloudFormation", "IaC deployment", 13.5),
]
for val, label, x in metrics:
    txt(x, 7.98, 2.5, 0.2, val, 10, ORANGE, True)
    txt(x, 8.15, 2.5, 0.2, label, 8, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide21_governance_deep_dive_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
