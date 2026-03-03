from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)
BG = RGBColor(13, 17, 33)
PLACEHOLDER_BG = RGBColor(18, 24, 45)

def make_slide():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return prs, slide

def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(slide, l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
        if i > 0: p.space_before = Pt(2)

def rect(slide, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def hline(slide, x, y, w, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def vline(slide, x, y, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def footer(slide):
    txt(slide, 0.6, 8.6, 12, 0.3,
        "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

# =====================================================================
# VARIANT A: Two clean columns — no boxes, accent lines only
# =====================================================================
prs_a, sl_a = make_slide()

txt(sl_a, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — MODEL & INFERENCE", 34, ORANGE, True)
txt(sl_a, 0.6, 0.65, 14, 0.3,
    "SageMaker Random Cut Forest training, batch inference, real-time API, and anomaly thresholding", 16, LIGHT)

# LEFT COLUMN
vline(sl_a, 0.6, 1.2, 6.3, PURPLE)

txt(sl_a, 0.9, 1.2, 6, 0.25, "TRAINING DATASET", 13, PURPLE, True)
hline(sl_a, 0.9, 1.47, 2.2, PURPLE)
mtxt(sl_a, 0.9, 1.55, 6.5, 1.2, [
    ("Raw features from Redshift telemetry:", WHITE, False),
    ("pressure (PSI) · temperature (°F) · vehicle_id · tire_id · timestamp", LIGHT, False),
    ("", LIGHT, False),
    ("Engineered features (ML ETL Glue job):", WHITE, False),
    ("delta_pressure (lag diff) · delta_temp (lag diff)", LIGHT, False),
    ("Normalization: z-score via Welford's online algorithm, stats persisted in SSM", MUTED, False),
], 10)

txt(sl_a, 0.9, 3.0, 6, 0.25, "MODEL CONFIGURATION", 13, PURPLE, True)
hline(sl_a, 0.9, 3.27, 2.8, PURPLE)
mtxt(sl_a, 0.9, 3.35, 6.5, 1.5, [
    ("Algorithm: SageMaker Random Cut Forest", SKY, True),
    ("Unsupervised anomaly detection — no labeled data required", LIGHT, False),
    ("Outputs anomaly score (0.0 – 1.0) per observation", LIGHT, False),
    ("", LIGHT, False),
    ("Hyperparameters:", WHITE, False),
    ("feature_dim · num_samples_per_tree: 256 · num_trees: 100", LIGHT, False),
    ("", LIGHT, False),
    ("Infrastructure:", WHITE, False),
    ("Training: 4× ml.m5.12xlarge · 400 GiB volume · max 1 hour", LIGHT, False),
    ("Model artifacts → S3 → SageMaker Model → Endpoint", MUTED, False),
], 10)

txt(sl_a, 0.9, 5.15, 6, 0.25, "TRAINING PIPELINE", 13, PURPLE, True)
hline(sl_a, 0.9, 5.42, 2.3, PURPLE)
mtxt(sl_a, 0.9, 5.5, 6.5, 1.5, [
    ("Step Function orchestration:", WHITE, False),
    ("1. ML ETL cleaner Lambda (prep training data)", LIGHT, False),
    ("2. Generate unique job ID", LIGHT, False),
    ("3. SageMaker CreateTrainingJob (RCF container from ECR)", LIGHT, False),
    ("4. CreateModel from trained artifacts", LIGHT, False),
    ("5. Update model name in SSM Parameter Store", LIGHT, False),
    ("6. CreateEndpointConfig → UpdateEndpoint (or Create if first run)", LIGHT, False),
], 10)

# RIGHT COLUMN
vline(sl_a, 8.3, 1.2, 6.3, GREEN)

txt(sl_a, 8.6, 1.2, 6, 0.25, "BATCH INFERENCE", 13, GREEN, True)
hline(sl_a, 8.6, 1.47, 2.0, GREEN)
mtxt(sl_a, 8.6, 1.55, 6.5, 1.5, [
    ("Daily Step Function:", WHITE, False),
    ("1. Lambda reads model name from SSM", LIGHT, False),
    ("2. SageMaker Batch Transform (ml.m5.12xlarge)", LIGHT, False),
    ("   InputFilter: features only ($[4:7])", MUTED, False),
    ("   OutputFilter: metadata + anomaly score", MUTED, False),
    ("3. Status polling Lambda (30s intervals)", LIGHT, False),
    ("4. Raw predictions → S3 (.csv.out)", LIGHT, False),
], 10)

txt(sl_a, 8.6, 3.3, 6, 0.25, "ANOMALY THRESHOLDING", 13, GREEN, True)
hline(sl_a, 8.6, 3.57, 2.8, GREEN)
mtxt(sl_a, 8.6, 3.65, 6.5, 0.8, [
    ("99th percentile of batch anomaly scores", YELLOW, True),
    ("Threshold updated in SSM each batch run", LIGHT, False),
    ("Score > threshold → anomaly alert triggered", LIGHT, False),
    ("Dynamic: adapts as fleet data distribution changes", MUTED, False),
], 10)

txt(sl_a, 8.6, 4.7, 6, 0.25, "REAL-TIME API", 13, GREEN, True)
hline(sl_a, 8.6, 4.97, 1.7, GREEN)
mtxt(sl_a, 8.6, 5.05, 6.5, 1.5, [
    ("POST /predict  (API Gateway + Lambda + SageMaker endpoint)", GREEN, True),
    ("", LIGHT, False),
    ("Request: { vehicle_id, tire_id, pressure, temperature,", LIGHT, False),
    ("           delta_pressure, delta_temp }", LIGHT, False),
    ("", LIGHT, False),
    ("Response: { anomaly_score: 0.754, prediction: true }", LIGHT, False),
    ("", LIGHT, False),
    ("Lambda normalizes inputs using SSM stats → invokes endpoint (CSV)", LIGHT, False),
    ("Rate: 100 req/s · Burst: 200 · Daily: 10K · API Key auth", MUTED, False),
], 10)

# Bottom
hline(sl_a, 0.6, 7.75, 14.8, ORANGE)
txt(sl_a, 0.6, 7.85, 14.8, 0.3,
    "Training: 4 features · RCF unsupervised · 4× ml.m5.12xlarge  |  Inference: daily batch + real-time API  |  Threshold: 99th pctl dynamic",
    10, ORANGE)

footer(sl_a)
prs_a.save(os.path.join(OUTDIR, "slide20b_model_varA.pptx"))
print("Saved: slide20b_model_varA.pptx")

# =====================================================================
# VARIANT B: Top model card + bottom clean text
# =====================================================================
prs_b, sl_b = make_slide()

txt(sl_b, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — MODEL & INFERENCE", 34, ORANGE, True)
txt(sl_b, 0.6, 0.65, 14, 0.3,
    "SageMaker Random Cut Forest training, batch inference, real-time API, and anomaly thresholding", 16, LIGHT)

# Single accent card for the model — the ONE box on the slide
rect(sl_b, 0.6, 1.15, 14.8, 1.8, CARD, SKY, 2)
txt(sl_b, 0.9, 1.2, 6, 0.3, "SageMaker Random Cut Forest", 18, SKY, True)
txt(sl_b, 0.9, 1.5, 14, 0.25, "Unsupervised anomaly detection — learns normal fleet patterns, scores deviations without labeled data", 12, LIGHT)

mtxt(sl_b, 0.9, 1.85, 4.5, 1.0, [
    ("Training: 4× ml.m5.12xlarge", WHITE, False),
    ("400 GiB volume · max 1 hour", LIGHT, False),
    ("RCF container from ECR", MUTED, False),
], 10)
mtxt(sl_b, 5.5, 1.85, 4.5, 1.0, [
    ("Hyperparameters:", WHITE, False),
    ("feature_dim · num_samples_per_tree: 256", LIGHT, False),
    ("num_trees: 100", MUTED, False),
], 10)
mtxt(sl_b, 10.2, 1.85, 5, 1.0, [
    ("Endpoint: ml.m5.xlarge (real-time)", WHITE, False),
    ("Batch: ml.m5.12xlarge (daily)", LIGHT, False),
    ("Threshold: 99th pctl → SSM", MUTED, False),
], 10)

# Below: three clean text sections with accent lines
# Dataset
txt(sl_b, 0.6, 3.2, 5, 0.25, "DATASET & FEATURES", 13, PURPLE, True)
hline(sl_b, 0.6, 3.47, 2.5, PURPLE)
mtxt(sl_b, 0.6, 3.55, 4.8, 1.8, [
    ("Source: Redshift Datashare (MMT telemetry)", WHITE, False),
    ("Hourly ingestion → Root ETL → ML ETL", LIGHT, False),
    ("", LIGHT, False),
    ("Raw: pressure, temperature, vehicle_id, tire_id, timestamp", LIGHT, False),
    ("Engineered: delta_pressure, delta_temp (lag diffs)", LIGHT, False),
    ("Normalized: z-score via Welford's online algorithm", LIGHT, False),
    ("Running stats persisted in SSM Parameter Store", MUTED, False),
    ("", LIGHT, False),
    ("Training cols: [pressure, temp, Δp, Δt]", PURPLE, False),
    ("Inference cols: [metadata] + training cols", MUTED, False),
], 10)

# Training Pipeline
txt(sl_b, 5.7, 3.2, 5, 0.25, "TRAINING PIPELINE", 13, GREEN, True)
hline(sl_b, 5.7, 3.47, 2.3, GREEN)
mtxt(sl_b, 5.7, 3.55, 4.8, 1.8, [
    ("Step Function orchestration:", WHITE, False),
    ("1. ML ETL cleaner Lambda", LIGHT, False),
    ("2. Generate unique training job ID", LIGHT, False),
    ("3. SageMaker CreateTrainingJob", LIGHT, False),
    ("4. Wait for completion (.sync)", LIGHT, False),
    ("5. CreateModel from artifacts", LIGHT, False),
    ("6. Update model name → SSM", LIGHT, False),
    ("7. CreateEndpointConfig", LIGHT, False),
    ("8. UpdateEndpoint (or Create)", LIGHT, False),
], 10)

# Inference + API
txt(sl_b, 10.8, 3.2, 5, 0.25, "INFERENCE & API", 13, ORANGE, True)
hline(sl_b, 10.8, 3.47, 2.0, ORANGE)
mtxt(sl_b, 10.8, 3.55, 4.6, 1.8, [
    ("Batch (daily Step Function):", WHITE, False),
    ("Lambda → Batch Transform", LIGHT, False),
    ("InputFilter: $[4:7] (features only)", MUTED, False),
    ("OutputFilter: metadata + score", MUTED, False),
    ("Polling: 30s intervals until done", LIGHT, False),
    ("", LIGHT, False),
    ("Real-time: POST /predict", GREEN, True),
    ("API Gateway + Lambda + endpoint", LIGHT, False),
    ("Normalize → invoke → score > threshold?", LIGHT, False),
    ("100 req/s · API Key auth · CORS", MUTED, False),
], 10)

# Why dual-path — clean text, no box
hline(sl_b, 0.6, 5.65, 14.8, RGBColor(35, 45, 72))
txt(sl_b, 0.6, 5.8, 14, 0.25, "WHY DUAL-PATH VALIDATION?", 12, ORANGE, True)
txt(sl_b, 0.6, 6.1, 7, 0.25,
    "ML path catches novel, unseen failure patterns across the fleet — no domain rules needed", 10, PURPLE)
txt(sl_b, 8.0, 6.1, 7, 0.25,
    "Filter path provides deterministic, explainable leak rates with severity classification", 10, YELLOW)

# Bottom
hline(sl_b, 0.6, 7.75, 14.8, ORANGE)
txt(sl_b, 0.6, 7.85, 14.8, 0.3,
    "4 features  ·  RCF unsupervised  ·  4× ml.m5.12xlarge  ·  Daily batch + real-time API  ·  99th pctl threshold  ·  Welford's normalization",
    10, ORANGE)

footer(sl_b)
prs_b.save(os.path.join(OUTDIR, "slide20b_model_varB.pptx"))
print("Saved: slide20b_model_varB.pptx")

# =====================================================================
# VARIANT C: Three equal columns — Dataset | Model | Inference
# =====================================================================
prs_c, sl_c = make_slide()

txt(sl_c, 0.6, 0.2, 14, 0.5, "PREDICTIVE MAINTENANCE — MODEL & INFERENCE", 34, ORANGE, True)
txt(sl_c, 0.6, 0.65, 14, 0.3,
    "SageMaker Random Cut Forest training, batch inference, real-time API, and anomaly thresholding", 16, LIGHT)

# Three columns with vertical accent lines
col_w = 4.7
col1 = 0.6
col2 = 5.6
col3 = 10.6

# Column 1: Dataset
vline(sl_c, col1, 1.15, 6.3, PURPLE)
txt(sl_c, col1 + 0.3, 1.15, col_w, 0.3, "DATASET & FEATURES", 14, PURPLE, True)
hline(sl_c, col1 + 0.3, 1.47, 2.5, PURPLE)

mtxt(sl_c, col1 + 0.3, 1.6, col_w - 0.3, 5.5, [
    ("Source", WHITE, True),
    ("Redshift Datashare — MMT tire telemetry", LIGHT, False),
    ("Hourly Lambda query → S3 raw bucket", LIGHT, False),
    ("", LIGHT, False),
    ("Root ETL (Glue Spark, hourly)", WHITE, True),
    ("Clean · convert units · weave tables", LIGHT, False),
    ("→ S3 partitioned by date/hour", MUTED, False),
    ("", LIGHT, False),
    ("ML ETL (daily Step Function)", WHITE, True),
    ("Resample to 1-day intervals", LIGHT, False),
    ("Group by vehicle_id + tire_id", LIGHT, False),
    ("", LIGHT, False),
    ("Raw Features", WHITE, True),
    ("pressure (PSI)", LIGHT, False),
    ("temperature (°F)", LIGHT, False),
    ("vehicle_id, tire_id, timestamp", MUTED, False),
    ("", LIGHT, False),
    ("Engineered Features", WHITE, True),
    ("delta_pressure (lag difference)", LIGHT, False),
    ("delta_temp (lag difference)", LIGHT, False),
    ("", LIGHT, False),
    ("Normalization", WHITE, True),
    ("Z-score: (val − mean) / std", LIGHT, False),
    ("Welford's online algorithm", LIGHT, False),
    ("Running stats → SSM", MUTED, False),
], 10)

# Column 2: Model
vline(sl_c, col2, 1.15, 6.3, SKY)
txt(sl_c, col2 + 0.3, 1.15, col_w, 0.3, "MODEL & TRAINING", 14, SKY, True)
hline(sl_c, col2 + 0.3, 1.47, 2.3, SKY)

mtxt(sl_c, col2 + 0.3, 1.6, col_w - 0.3, 5.5, [
    ("Algorithm", WHITE, True),
    ("SageMaker Random Cut Forest", SKY, True),
    ("Unsupervised anomaly detection", LIGHT, False),
    ("No labeled data required", LIGHT, False),
    ("Score output: 0.0 – 1.0", LIGHT, False),
    ("", LIGHT, False),
    ("Hyperparameters", WHITE, True),
    ("feature_dim (4 features)", LIGHT, False),
    ("num_samples_per_tree: 256", LIGHT, False),
    ("num_trees: 100", LIGHT, False),
    ("", LIGHT, False),
    ("Infrastructure", WHITE, True),
    ("Training: 4× ml.m5.12xlarge", LIGHT, False),
    ("Volume: 400 GiB per instance", LIGHT, False),
    ("Max runtime: 3600s (1 hour)", MUTED, False),
    ("", LIGHT, False),
    ("Pipeline (Step Function)", WHITE, True),
    ("1. ETL cleaner Lambda", LIGHT, False),
    ("2. CreateTrainingJob", LIGHT, False),
    ("3. CreateModel", LIGHT, False),
    ("4. Update SSM model name", LIGHT, False),
    ("5. Create/Update Endpoint", LIGHT, False),
], 10)

# Column 3: Inference
vline(sl_c, col3, 1.15, 6.3, GREEN)
txt(sl_c, col3 + 0.3, 1.15, col_w, 0.3, "INFERENCE & API", 14, GREEN, True)
hline(sl_c, col3 + 0.3, 1.47, 2.0, GREEN)

mtxt(sl_c, col3 + 0.3, 1.6, col_w - 0.3, 5.5, [
    ("Batch (daily Step Function)", WHITE, True),
    ("Lambda reads model from SSM", LIGHT, False),
    ("Batch Transform: ml.m5.12xlarge", LIGHT, False),
    ("InputFilter: $[4:7] (features)", MUTED, False),
    ("OutputFilter: metadata + score", MUTED, False),
    ("Polling: 30s until complete", LIGHT, False),
    ("→ S3 raw predictions (.csv.out)", LIGHT, False),
    ("", LIGHT, False),
    ("Thresholding", WHITE, True),
    ("99th percentile of batch scores", YELLOW, True),
    ("Updated in SSM each run", LIGHT, False),
    ("Dynamic: adapts to fleet changes", MUTED, False),
    ("", LIGHT, False),
    ("Real-time API", WHITE, True),
    ("POST /predict", GREEN, True),
    ("API Gateway + Lambda + endpoint", LIGHT, False),
    ("Input: 6 fields (vehicle, tire, sensors)", LIGHT, False),
    ("Output: { score, prediction }", LIGHT, False),
    ("Normalize → invoke → threshold", MUTED, False),
    ("", LIGHT, False),
    ("Rate Limits", WHITE, True),
    ("100 req/s · Burst: 200", LIGHT, False),
    ("Daily: 10K · API Key auth", MUTED, False),
], 10)

# Bottom
hline(sl_c, 0.6, 7.75, 14.8, ORANGE)
txt(sl_c, 0.6, 7.85, 14.8, 0.3,
    "4 features  ·  RCF unsupervised  ·  4× ml.m5.12xlarge  ·  Daily batch + real-time  ·  99th pctl threshold  ·  Welford's normalization",
    10, ORANGE)

footer(sl_c)
prs_c.save(os.path.join(OUTDIR, "slide20b_model_varC.pptx"))
print("Saved: slide20b_model_varC.pptx")
