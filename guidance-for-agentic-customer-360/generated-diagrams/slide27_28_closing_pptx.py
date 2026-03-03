from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
BG = RGBColor(13, 17, 33)

def make_slide():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return prs, slide

def txt(sl, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(sl, l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
        if i > 0: p.space_before = Pt(3)

def hline(sl, x, y, w, color):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def vline(sl, x, y, h, color):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def rect(sl, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def footer(sl):
    txt(sl, 0.6, 8.6, 12, 0.3,
        "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

# =====================================================================
# SLIDE: GETTING STARTED
# =====================================================================
prs_gs, sl_gs = make_slide()

txt(sl_gs, 0.6, 0.2, 14, 0.5, "GETTING STARTED", 34, ORANGE, True)
txt(sl_gs, 0.6, 0.65, 14, 0.3,
    "Three open-source solutions you can deploy today — MIT-0 licensed, production-ready", 16, LIGHT)

# Three columns: one per solution
col1, col2, col3 = 0.6, 5.6, 10.6

# --- Customer 360 ---
vline(sl_gs, col1, 1.3, 4.8, SKY)
txt(sl_gs, col1 + 0.3, 1.3, 4.5, 0.3, "CUSTOMER 360 ANALYTICS", 14, SKY, True)
hline(sl_gs, col1 + 0.3, 1.6, 3.2, SKY)
mtxt(sl_gs, col1 + 0.3, 1.7, 4.3, 4.3, [
    ("What you get:", WHITE, True),
    ("S3 data lake · Glue catalog (11 tables)", LIGHT, False),
    ("Athena views · QuickSight dashboards", LIGHT, False),
    ("Aurora pgvector · Bedrock Agent", LIGHT, False),
    ("500K synthetic customers included", LIGHT, False),
    ("", LIGHT, False),
    ("Deploy:", WHITE, True),
    ("6 CDK stacks · ~30 min to deploy", LIGHT, False),
    ("npm install → cdk deploy --all", MUTED, False),
    ("", LIGHT, False),
    ("Prerequisites:", WHITE, True),
    ("AWS account · CDK CLI · Node.js", LIGHT, False),
    ("QuickSight subscription (for dashboards)", MUTED, False),
    ("", LIGHT, False),
    ("Cost: ~$114/mo (serverless baseline)", GREEN, True),
], 10)

# --- Predictive Maintenance ---
vline(sl_gs, col2, 1.3, 4.8, PURPLE)
txt(sl_gs, col2 + 0.3, 1.3, 4.5, 0.3, "PREDICTIVE MAINTENANCE", 14, PURPLE, True)
hline(sl_gs, col2 + 0.3, 1.6, 3.0, PURPLE)
mtxt(sl_gs, col2 + 0.3, 1.7, 4.3, 4.3, [
    ("What you get:", WHITE, True),
    ("Dual-path prediction pipeline", LIGHT, False),
    ("SageMaker RCF training + inference", LIGHT, False),
    ("Physics-based slow leak detection", LIGHT, False),
    ("Real-time API + batch processing", LIGHT, False),
    ("Alerts system (DynamoDB + SQS)", LIGHT, False),
    ("", LIGHT, False),
    ("Deploy:", WHITE, True),
    ("CDK (Python) · make build && make deploy", LIGHT, False),
    ("Requires Redshift datashare setup", MUTED, False),
    ("", LIGHT, False),
    ("Prerequisites:", WHITE, True),
    ("AWS account · Python 3.12 · Poetry", LIGHT, False),
    ("CDK CLI · Redshift data source", MUTED, False),
    ("", LIGHT, False),
    ("Built with Amazon Middle Mile fleet", GREEN, True),
], 10)

# --- Data Mesh Foundation ---
vline(sl_gs, col3, 1.3, 4.8, GREEN)
txt(sl_gs, col3 + 0.3, 1.3, 4.5, 0.3, "DATA MESH FOUNDATION", 14, GREEN, True)
hline(sl_gs, col3 + 0.3, 1.6, 2.8, GREEN)
mtxt(sl_gs, col3 + 0.3, 1.7, 4.3, 4.3, [
    ("What you get:", WHITE, True),
    ("SageMaker Unified Studio domain", LIGHT, False),
    ("DataZone V2 with SSO", LIGHT, False),
    ("5 blueprints (ML, Glue, Redshift,", LIGHT, False),
    ("EMR, Bedrock)", LIGHT, False),
    ("VPC + networking + Glue catalog", LIGHT, False),
    ("", LIGHT, False),
    ("Deploy:", WHITE, True),
    ("CloudFormation · deploy-complete-", LIGHT, False),
    ("platform.sh (automated)", MUTED, False),
    ("", LIGHT, False),
    ("Prerequisites:", WHITE, True),
    ("AWS account · IAM Identity Center", LIGHT, False),
    ("AWS CLI · CloudFormation access", MUTED, False),
    ("", LIGHT, False),
    ("Connects all solutions together", GREEN, True),
], 10)

# Bottom bar
hline(sl_gs, 0.6, 6.35, 14.8, RGBColor(35, 45, 72))

txt(sl_gs, 0.6, 6.6, 14, 0.3, "ALL THREE SOLUTIONS", 14, ORANGE, True)
hline(sl_gs, 0.6, 6.9, 2.8, ORANGE)

mtxt(sl_gs, 0.6, 7.0, 7, 0.8, [
    ("License: MIT-0 (no attribution required)", WHITE, False),
    ("Language: TypeScript (CDK) + Python (Lambda/Glue)", LIGHT, False),
    ("Region: us-east-1 (default, configurable)", MUTED, False),
], 11)

mtxt(sl_gs, 8.0, 7.0, 7, 0.8, [
    ("Deploy independently or together", WHITE, False),
    ("Data Mesh Foundation unifies catalog across solutions", LIGHT, False),
    ("Modular — adopt one, two, or all three", MUTED, False),
], 11)

footer(sl_gs)
prs_gs.save(os.path.join(OUTDIR, "slide27_getting_started.pptx"))
print("Saved: slide27_getting_started.pptx")

# =====================================================================
# SLIDE: CALL TO ACTION / NEXT STEPS
# =====================================================================
prs_cta, sl_cta = make_slide()

txt(sl_cta, 0.6, 0.2, 14, 0.5, "NEXT STEPS", 34, ORANGE, True)
hline(sl_cta, 0.6, 0.7, 14.8, RGBColor(35, 45, 72))

# Three engagement tiers — clean text, numbered
txt(sl_cta, 0.6, 1.1, 14, 0.4, "1", 48, ORANGE, True)
txt(sl_cta, 1.4, 1.15, 13, 0.3, "EXPLORE", 20, WHITE, True)
txt(sl_cta, 1.4, 1.5, 13, 0.3,
    "Clone the repo and deploy to a sandbox account — see it running in under an hour", 14, LIGHT)
hline(sl_cta, 0.6, 2.05, 14.8, RGBColor(35, 45, 72))

txt(sl_cta, 0.6, 2.3, 14, 0.4, "2", 48, ORANGE, True)
txt(sl_cta, 1.4, 2.35, 13, 0.3, "WORKSHOP", 20, WHITE, True)
txt(sl_cta, 1.4, 2.7, 13, 0.3,
    "Half-day hands-on session — deploy the platform, connect your data, build your first dashboard and agent", 14, LIGHT)
hline(sl_cta, 0.6, 3.25, 14.8, RGBColor(35, 45, 72))

txt(sl_cta, 0.6, 3.5, 14, 0.4, "3", 48, ORANGE, True)
txt(sl_cta, 1.4, 3.55, 13, 0.3, "PILOT", 20, WHITE, True)
txt(sl_cta, 1.4, 3.9, 13, 0.3,
    "4-6 week engagement — adapt the platform to your OEM data, validate with real telemetry, measure outcomes", 14, LIGHT)
hline(sl_cta, 0.6, 4.45, 14.8, RGBColor(35, 45, 72))

# What you walk away with — per tier
txt(sl_cta, 0.6, 4.8, 14, 0.3, "WHAT YOU WALK AWAY WITH", 14, ORANGE, True)
hline(sl_cta, 0.6, 5.12, 3.5, ORANGE)

col1, col2, col3 = 0.6, 5.6, 10.6

vline(sl_cta, col1, 5.3, 2.0, SKY)
txt(sl_cta, col1 + 0.3, 5.3, 4.5, 0.25, "EXPLORE", 12, SKY, True)
mtxt(sl_cta, col1 + 0.3, 5.6, 4.3, 1.6, [
    ("Running platform in your account", WHITE, False),
    ("500K customer dataset loaded", LIGHT, False),
    ("QuickSight dashboards live", LIGHT, False),
    ("Bedrock Agent answering questions", LIGHT, False),
    ("Predictive pipeline processing", MUTED, False),
], 11)

vline(sl_cta, col2, 5.3, 2.0, PURPLE)
txt(sl_cta, col2 + 0.3, 5.3, 4.5, 0.25, "WORKSHOP", 12, PURPLE, True)
mtxt(sl_cta, col2 + 0.3, 5.6, 4.3, 1.6, [
    ("Everything in Explore, plus:", WHITE, False),
    ("Hands-on with your team", LIGHT, False),
    ("Custom agent prompts for your use case", LIGHT, False),
    ("Architecture review for production", LIGHT, False),
    ("Identified gaps and next steps", MUTED, False),
], 11)

vline(sl_cta, col3, 5.3, 2.0, GREEN)
txt(sl_cta, col3 + 0.3, 5.3, 4.5, 0.25, "PILOT", 12, GREEN, True)
mtxt(sl_cta, col3 + 0.3, 5.6, 4.3, 1.6, [
    ("Everything in Workshop, plus:", WHITE, False),
    ("Connected to your real data sources", LIGHT, False),
    ("Tuned ML models on your fleet", LIGHT, False),
    ("Production-ready architecture", LIGHT, False),
    ("Measured ROI and business case", MUTED, False),
], 11)

# Contact
hline(sl_cta, 0.6, 7.6, 14.8, ORANGE)
txt(sl_cta, 0.6, 7.75, 14, 0.4,
    "Let's start the conversation", 20, ORANGE, True, PP_ALIGN.CENTER)
txt(sl_cta, 0.6, 8.15, 14, 0.3,
    "Your AWS account team can schedule any of these engagements — or reach out directly",
    13, LIGHT, False, PP_ALIGN.CENTER)

footer(sl_cta)
prs_cta.save(os.path.join(OUTDIR, "slide28_next_steps.pptx"))
print("Saved: slide28_next_steps.pptx")
