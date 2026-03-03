from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=11, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.6, "PREDICTIVE MAINTENANCE DEEP DIVE", 36, ORANGE, True)
txt(0.6, 0.75, 14, 0.4,
    "ML-based tire failure prediction with 7-14 day advance warning using dual-path validation", 18, LIGHT)

# === LEFT: Pipeline Architecture (vertical flow) ===

# --- IoT Ingestion ---
rect(0.6, 1.4, 8.7, 0.8, CARD, SKY, 1.5)
txt(0.9, 1.45, 3, 0.25, "IoT Core Ingestion", 13, SKY, True)
txt(4.0, 1.45, 5, 0.25, "MQTT protocol · Tire pressure · Tread depth · Temperature", 11, LIGHT)
txt(0.9, 1.8, 8, 0.25, "Vehicle TCU → IoT Core → IoT Rules Engine → Kinesis Data Streams", 10, MUTED)

arrow(4.95, 2.2, 4.95, 2.45, SKY, 1.5)

# --- Feature Engineering ---
rect(0.6, 2.45, 8.7, 0.8, CARD, GREEN, 1.5)
txt(0.9, 2.5, 3, 0.25, "Feature Engineering", 13, GREEN, True)
txt(4.0, 2.5, 5, 0.25, "Glue ETL · S3 Feature Store", 11, LIGHT)
mtxt(0.9, 2.8, 8, 0.4, [
    ("Rolling averages · Rate of change · Seasonal adjustments · Fleet baselines", MUTED, False),
], 10)

arrow(4.95, 3.25, 4.95, 3.5, GREEN, 1.5)

# --- Dual-Path Validation (two boxes side by side) ---
txt(0.6, 3.5, 8, 0.25, "DUAL-PATH VALIDATION", 12, ORANGE, True)

# ML path
rect(0.6, 3.8, 4.15, 1.4, CARD, PURPLE, 1.5)
txt(0.9, 3.85, 3.5, 0.25, "ML Anomaly Detection", 12, PURPLE, True)
mtxt(0.9, 4.15, 3.5, 1.0, [
    ("SageMaker Random Cut Forest", WHITE, False),
    ("• Unsupervised anomaly scoring", LIGHT, False),
    ("• Multi-variate time series", LIGHT, False),
    ("• Anomaly score threshold: 3σ", LIGHT, False),
    ("• Trained on 12 months fleet data", LIGHT, False),
], 10)

# Physics path
rect(5.15, 3.8, 4.15, 1.4, CARD, YELLOW, 1.5)
txt(5.45, 3.85, 3.5, 0.25, "Physics-Based Filters", 12, YELLOW, True)
mtxt(5.45, 4.15, 3.5, 1.0, [
    ("Statistical validation layer", WHITE, False),
    ("• Pressure decay rate > 0.5 psi/day", LIGHT, False),
    ("• Tread depth below 3/32\"", LIGHT, False),
    ("• Temperature anomaly correlation", LIGHT, False),
    ("• Eliminates sensor noise / false pos", LIGHT, False),
], 10)

# Merge arrow
arrow(4.75, 5.2, 4.75, 5.45, ORANGE, 2)
arrow(5.35, 5.2, 5.35, 5.45, ORANGE, 2)

# --- Output ---
rect(0.6, 5.45, 8.7, 0.7, CARD, ORANGE, 2)
txt(0.9, 5.5, 3, 0.25, "Validated Alert", 13, ORANGE, True)
txt(4.0, 5.5, 5, 0.25, "Both paths must agree → high-confidence prediction", 11, LIGHT)
txt(0.9, 5.8, 8, 0.25, "Real-time SageMaker endpoint + batch scoring pipeline via Step Functions", 10, MUTED)

# === RIGHT SIDE: Timeline + Outcomes ===
txt(10.0, 1.4, 5.5, 0.3, "PREDICTION TIMELINE", 13, ORANGE, True)

# Timeline visualization
days = [
    ("Day 0", "Normal operation\nAll readings nominal", GREEN, GREEN),
    ("Day 7-10", "Early anomaly detected\nML score crosses threshold", YELLOW, YELLOW),
    ("Day 10-12", "Physics filter confirms\nPressure decay validated", ORANGE, ORANGE),
    ("Day 12-14", "Alert triggered\nService appointment scheduled", PINK, PINK),
    ("Day 14+", "Without prediction:\nRoadside failure", RGBColor(248, 113, 113), RGBColor(248, 113, 113)),
]

dy = 1.8
for label, desc, color, border in days:
    rect(10.0, dy, 5.5, 0.65, CARD, border, 1.5)
    txt(10.3, dy + 0.05, 1.5, 0.25, label, 12, color, True)
    txt(11.8, dy + 0.05, 3.5, 0.55, desc, 10, LIGHT)
    dy += 0.75

# === RIGHT: Outcomes ===
txt(10.0, 5.7, 5.5, 0.3, "OUTCOMES", 13, GREEN, True)

rect(10.0, 6.05, 5.5, 1.5, CARD, GREEN, 1)
mtxt(10.3, 6.15, 5, 1.3, [
    ("7-14 day advance warning", GREEN, True),
    ("before tire failure occurs", LIGHT, False),
    ("", LIGHT, False),
    ("False positive reduction via dual-path", YELLOW, True),
    ("ML alone: ~15% FP → Dual-path: <3% FP", LIGHT, False),
    ("", LIGHT, False),
    ("Automated service scheduling", SKY, True),
    ("Alert → EventBridge → Connect outreach", LIGHT, False),
], 10)

# === BOTTOM: Key Metrics Bar ===
rect(0.6, 7.7, 14.8, 0.55, CARD, ORANGE, 1)
metrics = [
    ("7-14 days", "advance warning", 0.9),
    ("<3% FP rate", "dual-path validated", 3.5),
    ("Real-time", "inference endpoint", 6.1),
    ("Fleet-wide", "batch scoring", 8.7),
    ("Step Functions", "orchestration", 11.3),
    ("~$114/mo", "total platform cost", 13.5),
]
for val, label, x in metrics:
    txt(x, 7.75, 2, 0.2, val, 12, ORANGE, True)
    txt(x, 7.97, 2, 0.2, label, 9, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide20_predictive_maint_deep_dive.pptx")
prs.save(out)
print(f"Saved: {out}")
