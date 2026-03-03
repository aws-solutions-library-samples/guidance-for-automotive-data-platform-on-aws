from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def bullet_list(s, l, t, w, h, items, sz=13, color=LIGHT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = False
        p.space_after = Pt(6)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

# === HEADER ===
circ(slide, 0.5, 0.2, 0.55, YELLOW, None)
txt(slide, 0.5, 0.24, 0.55, 0.45, "3", 22, DARK, True, PP_ALIGN.CENTER)
txt(slide, 1.2, 0.2, 12, 0.6, "EVENT STREAMING", 36, ORANGE, True)
txt(slide, 1.2, 0.75, 12, 0.4, "Amazon MSK — Managed Streaming for Apache Kafka", 22, YELLOW, True)

# === LEFT COLUMN — What it does ===
lx = 0.6
txt(slide, lx, 1.4, 7, 0.35, "What Amazon MSK Does", 18, WHITE, True)

rect(slide, lx, 1.85, 7.0, 4.6, CARD, YELLOW, 1.5)

bullet_list(slide, lx + 0.25, 1.95, 6.5, 4.4, [
    ("Fully managed Apache Kafka that decouples vehicle data", WHITE),
    ("ingestion from processing — durable, replayable, ordered", LIGHT),
    ("", LIGHT),
    ("▸  Buffers billions of vehicle events per day with configurable", LIGHT),
    ("   retention (hours to indefinite) for replay and reprocessing", LIGHT),
    ("", LIGHT),
    ("▸  Fan-out to multiple consumers — same telemetry stream", LIGHT),
    ("   feeds analytics, ML, alerting, and third-party systems", LIGHT),
    ("", LIGHT),
    ("▸  MSK Serverless eliminates capacity planning — auto-scales", LIGHT),
    ("   throughput with no brokers to provision or manage", LIGHT),
    ("", LIGHT),
    ("▸  MSK Connect deploys Kafka connectors (S3, Redshift,", LIGHT),
    ("   OpenSearch, Snowflake) with zero custom code", LIGHT),
    ("", LIGHT),
    ("▸  Native integration with IoT Core Rules Engine — route", LIGHT),
    ("   MQTT messages directly into Kafka topics", LIGHT),
], 13, LIGHT)

# === RIGHT COLUMN — Why MSK ===
rx = 7.9
txt(slide, rx, 1.4, 7, 0.35, "Why MSK over Alternatives", 18, WHITE, True)

rect(slide, rx, 1.85, 7.5, 4.6, CARD, ORANGE, 1.5)

bullet_list(slide, rx + 0.25, 1.95, 7.0, 4.4, [
    ("vs. Confluent Cloud", WHITE),
    ("MSK Serverless is 40-60% lower cost at automotive scale.", LIGHT),
    ("No per-partition pricing — critical when thousands of vehicles", LIGHT),
    ("each publish to dedicated topics", LIGHT),
    ("", LIGHT),
    ("vs. Redpanda", WHITE),
    ("MSK offers deeper AWS integration (PrivateLink, IAM auth,", LIGHT),
    ("CloudWatch metrics, MSK Connect). Redpanda is Kafka-", LIGHT),
    ("compatible but lacks native AWS service integrations", LIGHT),
    ("", LIGHT),
    ("vs. Amazon Kinesis", WHITE),
    ("Kafka's consumer group model and topic-based routing are", LIGHT),
    ("better suited for vehicle platforms with diverse downstream", LIGHT),
    ("consumers. Existing Kafka expertise transfers directly", LIGHT),
    ("", LIGHT),
    ("Open standard — no lock-in", WHITE),
    ("Apache Kafka protocol means portability across MSK,", LIGHT),
    ("Confluent, Redpanda, or self-managed clusters", LIGHT),
], 13, LIGHT)

# === BOTTOM — Key capabilities ===
py = 6.7
txt(slide, 0.6, py, 4, 0.3, "Key Capabilities", 14, MUTED, True)

products = [
    ("MSK Serverless", YELLOW, "Auto-scaling Kafka — no brokers to manage"),
    ("MSK Connect", GREEN, "Managed connectors to S3, Redshift, OpenSearch"),
    ("Schema Registry", PURPLE, "Enforce data contracts across vehicle signals"),
]

px = 0.6
for name, color, desc in products:
    rect(slide, px, py + 0.35, 4.8, 0.7, CARD, color, 1)
    txt(slide, px + 0.2, py + 0.38, 2.5, 0.3, name, 14, WHITE, True)
    txt(slide, px + 0.2, py + 0.65, 4.4, 0.3, desc, 11, LIGHT)
    px += 5.0

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide12_msk.pptx")
prs.save(out)
print(f"Saved: {out}")
