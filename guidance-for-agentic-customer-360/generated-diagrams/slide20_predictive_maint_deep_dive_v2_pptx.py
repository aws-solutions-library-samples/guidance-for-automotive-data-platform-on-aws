from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.15, 14, 0.5, "PREDICTIVE MAINTENANCE DEEP DIVE", 34, ORANGE, True)
txt(0.6, 0.6, 14, 0.35,
    "Dual-path tire failure prediction: ML anomaly detection + physics-based slow leak filtering", 16, LIGHT)

# =====================================================================
# LEFT COLUMN: DATA PIPELINE (0.5 - 7.8)
# =====================================================================

# --- 1. Data Source ---
rect(0.5, 1.1, 7.3, 0.6, CARD, SKY, 1.5)
txt(0.8, 1.15, 3, 0.2, "① Data Source", 12, SKY, True)
txt(3.5, 1.15, 4, 0.2, "Redshift Datashare (MMT telemetry)", 10, LIGHT)
txt(0.8, 1.4, 6.5, 0.2, "Hourly SQL query via Lambda → raw tire data to S3", 9, MUTED)

arrow(4.15, 1.7, 4.15, 1.9, SKY, 1)

# --- 2. Root ETL ---
rect(0.5, 1.9, 7.3, 0.75, CARD, GREEN, 1.5)
txt(0.8, 1.95, 3, 0.2, "② Root ETL (Glue)", 12, GREEN, True)
txt(4.0, 1.95, 3.5, 0.2, "Hourly · Spark · CSV output", 10, LIGHT)
mtxt(0.8, 2.2, 6.5, 0.4, [
    ("Clean erroneous values · Unit conversions · Weave multi-table by vehicle_id + timestamp", LIGHT, False),
    ("Output: partitioned by date/hour → feeds both ML and Filtering paths", MUTED, False),
], 9)

# Split arrows to ML and Filtering
arrow(2.5, 2.65, 2.5, 2.9, GREEN, 1)
arrow(5.8, 2.65, 5.8, 2.9, GREEN, 1)

# --- 3a. ML ETL (left sub-column) ---
rect(0.5, 2.9, 3.5, 1.1, CARD, PURPLE, 1.5)
txt(0.7, 2.95, 3, 0.2, "③ ML ETL (Daily)", 11, PURPLE, True)
mtxt(0.7, 3.2, 3.1, 0.75, [
    ("Step Function → Glue Job:", LIGHT, False),
    ("• Resample to 1-day intervals", LIGHT, False),
    ("• Engineer: delta_pressure, delta_temp", LIGHT, False),
    ("• Normalize: z-score (Welford's)", LIGHT, False),
    ("• Stats persisted in SSM", MUTED, False),
], 9)

# --- 3b. Filtering Path (right sub-column) ---
rect(4.3, 2.9, 3.5, 1.1, CARD, YELLOW, 1.5)
txt(4.5, 2.95, 3, 0.2, "③ Slow Leak Filter (Nightly)", 11, YELLOW, True)
mtxt(4.5, 3.2, 3.1, 0.75, [
    ("Step Function → Lambda → Glue:", LIGHT, False),
    ("• Quantile filter (10th pctl rolling)", LIGHT, False),
    ("• Daily aggregation", LIGHT, False),
    ("• Leak rate = ΔPSI / Δdays", LIGHT, False),
    ("• Severity: H>3, M>2, L>1 PSI/day", MUTED, False),
], 9)

# Arrows down from both paths
arrow(2.5, 4.0, 2.5, 4.2, PURPLE, 1)
arrow(5.8, 4.0, 5.8, 4.2, YELLOW, 1)

# --- 4a. ML Training + Inference (left) ---
rect(0.5, 4.2, 3.5, 1.35, CARD, PURPLE, 1.5)
txt(0.7, 4.25, 3, 0.2, "④ ML Training & Inference", 11, PURPLE, True)
mtxt(0.7, 4.5, 3.1, 1.0, [
    ("Training (Step Function):", WHITE, True),
    ("• SageMaker Random Cut Forest", LIGHT, False),
    ("• 4× ml.m5.12xlarge instances", LIGHT, False),
    ("• Hyperparams: feature_dim, num_trees", MUTED, False),
    ("Inference (Daily Step Function):", WHITE, True),
    ("• Batch Transform → ml.m5.12xlarge", LIGHT, False),
    ("• Anomaly score per tire per day", LIGHT, False),
], 9)

# --- 4b. Filtering output (right) ---
rect(4.3, 4.2, 3.5, 1.35, CARD, YELLOW, 1.5)
txt(4.5, 4.25, 3, 0.2, "④ Leak Detection Output", 11, YELLOW, True)
mtxt(4.5, 4.5, 3.1, 1.0, [
    ("Per tire results:", WHITE, True),
    ("• Leak rate (PSI/day)", LIGHT, False),
    ("• Severity (HIGH/MED/LOW)", LIGHT, False),
    ("• Time to 80 PSI threshold", LIGHT, False),
    ("Metrics to CloudWatch:", WHITE, True),
    ("• num_slow_leak_detected", LIGHT, False),
    ("• num_aaids, num_tires processed", LIGHT, False),
], 9)

# Merge arrows into Alerts
arrow(2.5, 5.55, 4.15, 5.8, ORANGE, 1.5)
arrow(5.8, 5.55, 4.15, 5.8, ORANGE, 1.5)

# --- 5. Alerts System ---
rect(0.5, 5.8, 7.3, 0.85, CARD, ORANGE, 2)
txt(0.8, 5.85, 3, 0.2, "⑤ Alerts System", 12, ORANGE, True)
mtxt(0.8, 6.1, 6.8, 0.5, [
    ("ML predictions + Filtering results → S3 upload triggers alerts-transformer Lambda", LIGHT, False),
    ("→ DynamoDB (PENDING) → SQS queue → alerts-processor Lambda (PROCESSED) → downstream CSV", LIGHT, False),
    ("Failed alerts → DLQ → alerts-cleaner Lambda (FAILED) · 99th percentile anomaly threshold", MUTED, False),
], 9)

# =====================================================================
# RIGHT COLUMN: Training Details + API + Features (8.1 - 15.4)
# =====================================================================

# --- Training Dataset & Features ---
txt(8.1, 1.1, 7, 0.25, "TRAINING DATASET & FEATURES", 12, ORANGE, True)

rect(8.1, 1.4, 7.3, 1.65, CARD, PURPLE, 1.5)
txt(8.4, 1.45, 6.5, 0.2, "Data Pipeline: Redshift → Root ETL → ML ETL → SageMaker", 10, LIGHT)

mtxt(8.4, 1.7, 3.3, 1.3, [
    ("Raw Features (from telemetry):", WHITE, True),
    ("• pressure (tire PSI)", LIGHT, False),
    ("• temperature (°F)", LIGHT, False),
    ("• vehicle_id, tire_id", LIGHT, False),
    ("• timestamp", LIGHT, False),
], 9)

mtxt(11.8, 1.7, 3.3, 1.3, [
    ("Engineered Features:", WHITE, True),
    ("• delta_pressure (lag diff)", LIGHT, False),
    ("• delta_temp (lag diff)", LIGHT, False),
    ("Normalization:", WHITE, True),
    ("• Z-score with running stats", LIGHT, False),
    ("• Welford's online algorithm", MUTED, False),
], 9)

# --- Model Configuration ---
txt(8.1, 3.2, 7, 0.25, "MODEL CONFIGURATION", 12, ORANGE, True)

rect(8.1, 3.5, 7.3, 1.3, CARD, SKY, 1.5)
mtxt(8.4, 3.55, 3.3, 1.2, [
    ("Algorithm:", WHITE, True),
    ("SageMaker Random Cut Forest", SKY, True),
    ("• Unsupervised anomaly detection", LIGHT, False),
    ("• No labeled data required", LIGHT, False),
    ("• Anomaly score output (0.0–1.0)", LIGHT, False),
], 9)

mtxt(11.8, 3.55, 3.3, 1.2, [
    ("Infrastructure:", WHITE, True),
    ("Training: 4× ml.m5.12xlarge", SKY, False),
    ("Batch: 1× ml.m5.12xlarge", SKY, False),
    ("Endpoint: 1× ml.m5.xlarge", SKY, False),
    ("Threshold: 99th percentile", MUTED, False),
], 9)

# --- Real-Time API ---
txt(8.1, 5.0, 7, 0.25, "REAL-TIME PREDICTION API", 12, ORANGE, True)

rect(8.1, 5.3, 7.3, 1.1, CARD, GREEN, 1.5)
mtxt(8.4, 5.35, 3.3, 1.0, [
    ("POST /predict", GREEN, True),
    ("Input: vehicle_id, tire_id,", LIGHT, False),
    ("  pressure, temperature,", LIGHT, False),
    ("  delta_pressure, delta_temp", LIGHT, False),
], 9)

mtxt(11.8, 5.35, 3.3, 1.0, [
    ("Response:", GREEN, True),
    ("{ anomaly_score: 0.75,", LIGHT, False),
    ("  prediction: true }", LIGHT, False),
    ("API Gateway + Lambda + SageMaker", MUTED, False),
], 9)

# --- Dual-Path Validation Summary ---
txt(8.1, 6.6, 7, 0.25, "DUAL-PATH VALIDATION", 12, ORANGE, True)

rect(8.1, 6.9, 3.5, 0.75, CARD, PURPLE, 1)
mtxt(8.4, 6.95, 3, 0.65, [
    ("ML Path (Random Cut Forest)", PURPLE, True),
    ("Unsupervised anomaly scoring", LIGHT, False),
    ("Catches novel failure patterns", MUTED, False),
], 9)

rect(11.9, 6.9, 3.5, 0.75, CARD, YELLOW, 1)
mtxt(12.2, 6.95, 3, 0.65, [
    ("Filter Path (Physics-based)", YELLOW, True),
    ("Leak rate + severity classification", LIGHT, False),
    ("Deterministic, explainable results", MUTED, False),
], 9)

# === BOTTOM: Key Metrics Bar ===
rect(0.5, 7.9, 14.9, 0.45, CARD, ORANGE, 1)
metrics = [
    ("7-14 day", "advance warning", 0.8),
    ("4 features", "pressure, temp, Δp, Δt", 3.0),
    ("Dual-path", "ML + physics filters", 5.3),
    ("Real-time", "API + batch inference", 7.6),
    ("99th pctl", "anomaly threshold", 9.9),
    ("3 severity", "HIGH / MED / LOW", 12.2),
]
for val, label, x in metrics:
    txt(x, 7.93, 2, 0.2, val, 11, ORANGE, True)
    txt(x, 8.12, 2, 0.2, label, 8, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide20_predictive_maint_deep_dive_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
