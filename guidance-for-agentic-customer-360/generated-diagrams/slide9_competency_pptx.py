from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import math, os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(s, l, t, w, h, lines):
    """Multi-line textbox: lines = [(text, sz, color, bold, align), ...]"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, sz, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(4)

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def pill(s, l, t, text, color, sz=11):
    pw = len(text) * 0.085 + 0.3
    ph = 0.28
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(pw), Inches(ph))
    dark = RGBColor(color[0] // 6, color[1] // 6, color[2] // 6)
    sh.fill.solid(); sh.fill.fore_color.rgb = dark
    sh.line.color.rgb = color; sh.line.width = Pt(1)
    tb = sh.text_frame; tb.word_wrap = False
    p = tb.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz); p.font.color.rgb = color
    return pw

# === TITLE ===
txt(slide, 0.6, 0.2, 14, 0.6, "AWS AUTOMOTIVE COMPETENCY", 36, ORANGE, True)
txt(slide, 0.6, 0.75, 14, 0.4,
    "Validated partners with proven expertise across 8 automotive workloads", 20, LIGHT)

# === LEFT SIDE — What it is + Benefits ===

# What it is card
rect(slide, 0.6, 1.5, 7.2, 2.4, CARD, PURPLE, 1.5)
txt(slide, 0.9, 1.6, 6.5, 0.35, "What is the AWS Automotive Competency?", 18, PURPLE, True)
mtxt(slide, 0.9, 2.05, 6.5, 1.7, [
    ("A rigorous technical validation program that identifies AWS Partners", 13, LIGHT, False, PP_ALIGN.LEFT),
    ("with demonstrated expertise in automotive-specific cloud solutions.", 13, LIGHT, False, PP_ALIGN.LEFT),
    ("Partners undergo vetted customer reference reviews and technical", 13, LIGHT, False, PP_ALIGN.LEFT),
    ("assessments before earning the competency designation.", 13, LIGHT, False, PP_ALIGN.LEFT),
])

# Benefits for Partners — 3 cards
benefits = [
    ("Recognition", SKY,
     "Partner badge & priority placement in AWS Partner Solutions Finder. "
     "Improved co-sell score and direct access to AWS Sales teams."),
    ("Financial Incentives", GREEN,
     "Marketing Development Funds (MDF) for go-to-market activities. "
     "Migration Acceleration Program (MAP) funding for eligible workloads."),
    ("Access to AWS Experts", ORANGE,
     "AWS PartnerEquip enablement series, GTM Acceleration Workshops, "
     "partner-led webinars, and APN Customer Engagements (ACE) program."),
]

txt(slide, 0.6, 4.15, 7, 0.4, "Benefits for Partners", 20, WHITE, True)

bx = 0.6
bw = 2.3
bgap = 0.15
for i, (title, color, desc) in enumerate(benefits):
    x = bx + i * (bw + bgap)
    y = 4.6
    rect(slide, x, y, bw, 2.8, CARD, color, 1.5)

    # Icon circle
    circ(slide, x + bw/2 - 0.3, y + 0.2, 0.6, DARK, color, 2)
    txt(slide, x + bw/2 - 0.3, y + 0.28, 0.6, 0.45, str(i+1), 20, color, True, PP_ALIGN.CENTER)

    # Title
    txt(slide, x + 0.15, y + 0.95, bw - 0.3, 0.35, title, 15, WHITE, True, PP_ALIGN.CENTER)

    # Description
    txt(slide, x + 0.15, y + 1.35, bw - 0.3, 1.3, desc, 11, LIGHT, False, PP_ALIGN.LEFT)

# === RIGHT SIDE — 8 Workload Categories in a circle ===
txt(slide, 8.4, 1.5, 7, 0.4, "8 Automotive Workload Categories", 20, WHITE, True)

workloads = [
    ("Autonomous\nMobility", SKY),
    ("Software-Defined\nVehicle", GREEN),
    ("Connected\nMobility", YELLOW),
    ("Sustainability", ORANGE),
    ("Digital Customer\nExperience", PINK),
    ("Product\nEngineering", PURPLE),
    ("Manufacturing", RED),
    ("Supply\nChain", SKY),
]

# Circle layout
cx, cy = 11.8, 5.0
radius = 2.2
node_sz = 1.35

positions = []
for i in range(8):
    angle = math.radians(-90 + i * 45)
    x = cx + radius * math.cos(angle) - node_sz / 2
    y = cy + radius * math.sin(angle) - node_sz / 2
    positions.append((x, y))

# Dashed connectors
for i in range(8):
    ni = (i + 1) % 8
    x1 = positions[i][0] + node_sz / 2
    y1 = positions[i][1] + node_sz / 2
    x2 = positions[ni][0] + node_sz / 2
    y2 = positions[ni][1] + node_sz / 2
    dx, dy = x2 - x1, y2 - y1
    dist = math.sqrt(dx*dx + dy*dy)
    off = node_sz * 0.52
    sx, sy = x1 + dx/dist * off, y1 + dy/dist * off
    ex, ey = x2 - dx/dist * off, y2 - dy/dist * off
    conn = slide.shapes.add_connector(1, Inches(sx), Inches(sy), Inches(ex), Inches(ey))
    conn.line.color.rgb = RGBColor(35, 45, 72)
    conn.line.width = Pt(1)
    conn.line.dash_style = 2

# Draw workload nodes
for i, (name, color) in enumerate(workloads):
    x, y = positions[i]
    circ(slide, x, y, node_sz, CARD, color, 2)
    txt(slide, x + 0.05, y + node_sz/2 - 0.25, node_sz - 0.1, 0.55, name, 10, WHITE, True, PP_ALIGN.CENTER)

# Center label
circ(slide, cx - 0.7, cy - 0.7, 1.4, DARK, ORANGE, 2)
txt(slide, cx - 0.6, cy - 0.35, 1.2, 0.3, "AWS", 14, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.6, cy - 0.05, 1.2, 0.3, "Automotive", 12, ORANGE, True, PP_ALIGN.CENTER)
txt(slide, cx - 0.6, cy + 0.2, 1.2, 0.3, "Competency", 11, ORANGE, False, PP_ALIGN.CENTER)

# === BOTTOM — Launch partners bar ===
ly = 7.7
rect(slide, 0.6, ly, 14.8, 0.65, CARD)
txt(slide, 0.8, ly + 0.05, 3, 0.25, "Launch Partners include:", 12, MUTED, True)

launch_partners = [
    "WirelessCar", "Sibros", "Capgemini", "HCLTech", "Deloitte",
    "Accenture", "TCS", "Wipro", "Siemens", "MongoDB",
]
px = 4.0
for p in launch_partners:
    pw = pill(slide, px, ly + 0.18, p, PURPLE, 10)
    px += pw + 0.12

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide9_competency.pptx")
prs.save(out)
print(f"Saved: {out}")
