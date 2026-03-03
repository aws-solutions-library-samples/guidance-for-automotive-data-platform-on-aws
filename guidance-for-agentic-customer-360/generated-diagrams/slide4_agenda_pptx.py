from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
CARD = RGBColor(20, 28, 52)

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color=CARD, border_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# Title
add_text(slide, 0.6, 0.3, 8, 0.7, "TODAY'S AGENDA", 36, ORANGE, bold=True)
add_text(slide, 0.6, 0.8, 10, 0.5, "Vehicle Technology & Connected — Deep Dive", 20, LIGHT)

# Strategic domains bar
add_rect(slide, 0.6, 1.3, 14.8, 0.4, RGBColor(30, 38, 58))
add_text(slide, 0.8, 1.33, 2, 0.35, "Strategic Domains:", 13, MUTED)
add_text(slide, 3.0, 1.33, 2, 0.35, "Vehicle Technology", 13, SKY, bold=True)
add_text(slide, 5.0, 1.33, 1.5, 0.35, "Connected", 13, SKY, bold=True)
add_text(slide, 6.5, 1.33, 2.5, 0.35, "Product Engineering", 13, MUTED)
add_text(slide, 9.0, 1.33, 2.5, 0.35, "Smart Manufacturing", 13, MUTED)

# Sections
sections = [
    ("01", "Connected", "Modernizing the connected vehicle stack", SKY,
     "Cellular (ACS)  ·  Connectivity (IoT Core)  ·  Streaming (MSK)  ·  Processing (Flink)  ·  Observability"),
    ("02", "Automotive Data Platform", "Turning vehicle data into actionable insights", GREEN,
     "Data lake (S3)  ·  Analytics (Athena)  ·  ML & predictions (SageMaker)  ·  Data governance"),
    ("03", "Digital Customer Experience", "Engaging customers across every touchpoint", PINK,
     "Contact center (Connect)  ·  CRM integration  ·  Personalization  ·  After-sales"),
    ("04", "Fleet Management", "Optimizing fleet operations at scale", PURPLE,
     "Telematics & tracking  ·  Predictive maintenance  ·  Route optimization  ·  Driver safety"),
]

y_start = 2.0
card_h = 1.5
gap = 0.15

for i, (num, title, subtitle, color, topics) in enumerate(sections):
    y = y_start + i * (card_h + gap)

    # Card background
    add_rect(slide, 1.2, y, 13.8, card_h, CARD, color)

    # Number
    add_text(slide, 0.4, y + 0.4, 0.8, 0.6, num, 28, color, bold=True)

    # Title
    add_text(slide, 1.5, y + 0.15, 10, 0.5, title, 22, WHITE, bold=True)

    # Subtitle
    add_text(slide, 1.5, y + 0.55, 12, 0.4, subtitle, 14, LIGHT)

    # Topics
    add_text(slide, 1.5, y + 0.95, 13, 0.4, topics, 12, color)

# Footer
add_text(slide, 0.6, 8.6, 12, 0.3, "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide4_agenda.pptx")
prs.save(out)
print(f"Saved: {out}")
