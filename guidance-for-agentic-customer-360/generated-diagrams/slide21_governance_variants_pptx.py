from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
BG = RGBColor(13, 17, 33)
PLACEHOLDER_BG = RGBColor(18, 24, 45)

def make_slide():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return prs, slide

def txt(sl, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(sl, l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
        if i > 0: p.space_before = Pt(2)

def hline(sl, x, y, w, color):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def vline(sl, x, y, h, color):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def placeholder(sl, l, t, w, h, label="← INSERT ARCHITECTURE DIAGRAM →"):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PLACEHOLDER_BG
    sh.line.color.rgb = ORANGE; sh.line.width = Pt(1.5); sh.line.dash_style = 4
    txt(sl, l + 0.2, t + h/2 - 0.15, w - 0.4, 0.3, label, 12, ORANGE, False, PP_ALIGN.CENTER)

def rect(sl, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def footer(sl):
    txt(sl, 0.6, 8.6, 12, 0.3,
        "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

# =====================================================================
# VARIANT A: Left diagram + right three-section text
# =====================================================================
prs_a, sl_a = make_slide()

txt(sl_a, 0.6, 0.2, 14, 0.5, "DATA MESH FOUNDATION — ARCHITECTURE", 34, ORANGE, True)
txt(sl_a, 0.6, 0.65, 14, 0.3,
    "SageMaker Unified Studio (DataZone V2) — federated catalog, domain ownership, governed self-service", 16, LIGHT)

placeholder(sl_a, 0.6, 1.15, 8.5, 5.5,
    "← INSERT ARCHITECTURE DIAGRAM →")

# Right side: clean text with accent lines
x = 9.5

txt(sl_a, x, 1.15, 6, 0.25, "UNIFIED STUDIO DOMAIN", 13, ORANGE, True)
hline(sl_a, x, 1.42, 3.2, ORANGE)
mtxt(sl_a, x, 1.5, 5.8, 1.2, [
    ("Domain: automotive-data-platform", WHITE, False),
    ("DataZone V2 with IAM Identity Center (SSO)", LIGHT, False),
    ("Automatic user assignment · Portal URL for web access", LIGHT, False),
    ("DomainExecutionRole trusts: DataZone, SageMaker,", MUTED, False),
    ("Lake Formation, Glue — scoped to account resources", MUTED, False),
], 10)

txt(sl_a, x, 2.9, 6, 0.25, "BLUEPRINTS", 13, SKY, True)
hline(sl_a, x, 3.17, 1.5, SKY)
mtxt(sl_a, x, 3.25, 5.8, 1.2, [
    ("ML Analysis — SageMaker notebooks & training", LIGHT, False),
    ("Tooling — Glue ETL, Athena queries", LIGHT, False),
    ("Redshift — data warehouse analytics", LIGHT, False),
    ("EMR — big data processing", LIGHT, False),
    ("Bedrock — generative AI workloads", LIGHT, False),
], 10)

txt(sl_a, x, 4.65, 6, 0.25, "INFRASTRUCTURE", 13, GREEN, True)
hline(sl_a, x, 4.92, 2.0, GREEN)
mtxt(sl_a, x, 5.0, 5.8, 1.0, [
    ("VPC (10.0.0.0/16) · 3 private subnets across AZs", LIGHT, False),
    ("NAT Gateway · VPC endpoints (S3, Glue, SageMaker, Athena)", LIGHT, False),
    ("Encrypted S3 buckets (versioned, access-logged)", LIGHT, False),
    ("Glue database for federated catalog", MUTED, False),
], 10)

txt(sl_a, x, 6.2, 6, 0.25, "DATA PUBLISHING", 13, PURPLE, True)
hline(sl_a, x, 6.47, 2.0, PURPLE)
mtxt(sl_a, x, 6.55, 5.8, 0.8, [
    ("Glue databases → DataZone discovery job → accept assets → publish", LIGHT, False),
    ("Assets: Customer 360, vehicle telemetry, tire prediction, sales", LIGHT, False),
    ("Searchable catalog with lineage tracking", MUTED, False),
], 10)

hline(sl_a, 0.6, 7.75, 14.8, ORANGE)
txt(sl_a, 0.6, 7.85, 14.8, 0.3,
    "DataZone V2  ·  SSO (IAM Identity Center)  ·  5 blueprints  ·  3-AZ VPC  ·  Federated Glue catalog  ·  CloudFormation IaC",
    10, ORANGE)

footer(sl_a)
prs_a.save(os.path.join(OUTDIR, "slide21_governance_varA.pptx"))
print("Saved: slide21_governance_varA.pptx")

# =====================================================================
# VARIANT B: Top diagram + bottom three-column text
# =====================================================================
prs_b, sl_b = make_slide()

txt(sl_b, 0.6, 0.2, 14, 0.5, "DATA MESH FOUNDATION — ARCHITECTURE", 34, ORANGE, True)
txt(sl_b, 0.6, 0.65, 14, 0.3,
    "SageMaker Unified Studio (DataZone V2) — federated catalog, domain ownership, governed self-service", 16, LIGHT)

placeholder(sl_b, 0.6, 1.15, 14.8, 3.5,
    "← INSERT ARCHITECTURE DIAGRAM →")

# Three columns below
col1, col2, col3 = 0.6, 5.6, 10.6
col_w = 4.7

# Col 1: Platform
vline(sl_b, col1, 4.9, 2.6, ORANGE)
txt(sl_b, col1 + 0.3, 4.9, col_w, 0.25, "PLATFORM", 14, ORANGE, True)
hline(sl_b, col1 + 0.3, 5.17, 1.3, ORANGE)
mtxt(sl_b, col1 + 0.3, 5.25, col_w - 0.3, 2.2, [
    ("SageMaker Unified Studio", WHITE, True),
    ("DataZone V2 domain", LIGHT, False),
    ("IAM Identity Center (SSO)", LIGHT, False),
    ("Automatic user assignment", LIGHT, False),
    ("Web portal for all personas", MUTED, False),
    ("", LIGHT, False),
    ("CloudFormation deployment:", WHITE, True),
    ("shared-resources → datazone-domain", LIGHT, False),
    ("→ blueprint-roles → enable-blueprints", MUTED, False),
], 10)

# Col 2: Blueprints & Infra
vline(sl_b, col2, 4.9, 2.6, SKY)
txt(sl_b, col2 + 0.3, 4.9, col_w, 0.25, "BLUEPRINTS & INFRA", 14, SKY, True)
hline(sl_b, col2 + 0.3, 5.17, 2.5, SKY)
mtxt(sl_b, col2 + 0.3, 5.25, col_w - 0.3, 2.2, [
    ("5 enabled blueprints:", WHITE, True),
    ("ML Analysis (SageMaker)", LIGHT, False),
    ("Tooling (Glue, Athena)", LIGHT, False),
    ("Redshift · EMR · Bedrock", LIGHT, False),
    ("", LIGHT, False),
    ("Networking:", WHITE, True),
    ("VPC 10.0.0.0/16 · 3 AZ subnets", LIGHT, False),
    ("NAT Gateway · VPC endpoints", LIGHT, False),
    ("S3 encrypted, versioned, logged", MUTED, False),
], 10)

# Col 3: Catalog & Personas
vline(sl_b, col3, 4.9, 2.6, GREEN)
txt(sl_b, col3 + 0.3, 4.9, col_w, 0.25, "CATALOG & PERSONAS", 14, GREEN, True)
hline(sl_b, col3 + 0.3, 5.17, 2.5, GREEN)
mtxt(sl_b, col3 + 0.3, 5.25, col_w - 0.3, 2.2, [
    ("Data publishing:", WHITE, True),
    ("Glue → DataZone discovery → publish", LIGHT, False),
    ("Searchable catalog + lineage", LIGHT, False),
    ("", LIGHT, False),
    ("Data Engineer:", SKY, True),
    ("Glue ETL · Athena · publish datasets", LIGHT, False),
    ("ML Engineer:", PURPLE, True),
    ("SageMaker notebooks · model training", LIGHT, False),
    ("Analyst: QuickSight dashboards", MUTED, False),
], 10)

hline(sl_b, 0.6, 7.75, 14.8, ORANGE)
txt(sl_b, 0.6, 7.85, 14.8, 0.3,
    "DataZone V2  ·  SSO  ·  5 blueprints  ·  3-AZ VPC  ·  Federated catalog  ·  Domain ownership  ·  CloudFormation IaC",
    10, ORANGE)

footer(sl_b)
prs_b.save(os.path.join(OUTDIR, "slide21_governance_varB.pptx"))
print("Saved: slide21_governance_varB.pptx")

# =====================================================================
# VARIANT C: Single hero card + diagram placeholder + clean text
# =====================================================================
prs_c, sl_c = make_slide()

txt(sl_c, 0.6, 0.2, 14, 0.5, "DATA MESH FOUNDATION — ARCHITECTURE", 34, ORANGE, True)
txt(sl_c, 0.6, 0.65, 14, 0.3,
    "SageMaker Unified Studio (DataZone V2) — federated catalog, domain ownership, governed self-service", 16, LIGHT)

# One hero card — the domain
rect(sl_c, 0.6, 1.15, 14.8, 1.2, CARD, ORANGE, 2)
txt(sl_c, 0.9, 1.2, 5, 0.25, "SageMaker Unified Studio Domain", 16, ORANGE, True)
txt(sl_c, 0.9, 1.5, 14, 0.25,
    "automotive-data-platform  ·  DataZone V2  ·  IAM Identity Center SSO  ·  Automatic user assignment  ·  Web portal access", 11, LIGHT)
mtxt(sl_c, 0.9, 1.85, 14, 0.4, [
    ("Blueprints: ML Analysis · Tooling (Glue/Athena) · Redshift · EMR · Bedrock    |    "
     "Roles: DomainExecution · BlueprintProvisioning · DomainService", MUTED, False),
], 9)

# Diagram placeholder — middle
placeholder(sl_c, 0.6, 2.6, 14.8, 3.0,
    "← INSERT ARCHITECTURE DIAGRAM: Domain → Projects → Blueprints → Data Sources →")

# Below: two sections side by side, no boxes
txt(sl_c, 0.6, 5.85, 7, 0.25, "INFRASTRUCTURE", 13, GREEN, True)
hline(sl_c, 0.6, 6.12, 2.0, GREEN)
mtxt(sl_c, 0.6, 6.2, 7, 1.3, [
    ("Networking", WHITE, True),
    ("VPC (10.0.0.0/16) with 3 private subnets across AZs", LIGHT, False),
    ("NAT Gateway for outbound · VPC endpoints for S3, Glue, SageMaker, Athena, Redshift", LIGHT, False),
    ("", LIGHT, False),
    ("Storage", WHITE, True),
    ("DataZone environment bucket (encrypted, versioned, access-logged)", LIGHT, False),
    ("Glue database for federated catalog · Per-solution S3 data lakes", MUTED, False),
], 10)

txt(sl_c, 8.3, 5.85, 7, 0.25, "DATA CATALOG & PERSONAS", 13, PURPLE, True)
hline(sl_c, 8.3, 6.12, 3.0, PURPLE)
mtxt(sl_c, 8.3, 6.2, 7, 1.3, [
    ("Publishing pipeline", WHITE, True),
    ("Glue databases → DataZone create data source → run discovery → accept → publish", LIGHT, False),
    ("Assets: Customer 360, vehicle telemetry, tire prediction, sales, weather", MUTED, False),
    ("", LIGHT, False),
    ("Persona workflows", WHITE, True),
    ("Data Engineer: Glue ETL, Athena validation, publish to catalog", LIGHT, False),
    ("ML Engineer: SageMaker notebooks, model training, experiment tracking", LIGHT, False),
], 10)

hline(sl_c, 0.6, 7.75, 14.8, ORANGE)
txt(sl_c, 0.6, 7.85, 14.8, 0.3,
    "DataZone V2  ·  SSO  ·  5 blueprints  ·  3-AZ VPC  ·  Federated catalog  ·  Domain ownership  ·  CloudFormation IaC",
    10, ORANGE)

footer(sl_c)
prs_c.save(os.path.join(OUTDIR, "slide21_governance_varC.pptx"))
print("Saved: slide21_governance_varC.pptx")
