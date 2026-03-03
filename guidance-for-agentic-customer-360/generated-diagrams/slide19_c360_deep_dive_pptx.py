from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=11, color=LIGHT, align=PP_ALIGN.LEFT):
    """Multi-line text helper."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.6, "CUSTOMER 360 DEEP DIVE", 36, ORANGE, True)
txt(0.6, 0.75, 14, 0.4,
    "Unified customer profiles with AI-powered root cause analysis across 11 datasets", 18, LIGHT)

# === LEFT SIDE: Architecture Flow (vertical) ===

# --- Data Sources row ---
txt(0.6, 1.4, 3, 0.3, "DATA SOURCES", 13, SKY, True)
sources = ["CRM", "DMS", "Telemetry", "Service\nRecords", "Contact\nCenter", "Surveys"]
sx = 0.6
for name in sources:
    rect(sx, 1.75, 1.35, 0.7, CARD, SKY, 1.5)
    txt(sx + 0.05, 1.8, 1.25, 0.6, name, 10, WHITE, False, PP_ALIGN.CENTER)
    sx += 1.45

# Arrows down
for i in range(6):
    x = 0.6 + i * 1.45 + 0.675
    arrow(x, 2.45, x, 2.7, SKY, 1)

# --- S3 Data Lake ---
rect(0.6, 2.7, 8.7, 0.65, CARD, GREEN, 1.5)
txt(0.9, 2.75, 2.5, 0.25, "S3 Data Lake", 13, GREEN, True)
txt(3.5, 2.75, 5.5, 0.25, "Glue Crawlers → Glue Catalog → Athena (ad-hoc SQL)", 11, LIGHT)
txt(0.9, 3.05, 8, 0.25, "500K customers · 1.4M interactions · 900K service records · 11 datasets", 10, MUTED)

# Arrow down
arrow(4.95, 3.35, 4.95, 3.6, GREEN, 1.5)

# --- Processing: two paths ---
# Left path: QuickSight
rect(0.6, 3.6, 4.0, 1.2, CARD, YELLOW, 1.5)
txt(0.9, 3.65, 3.5, 0.25, "QuickSight Dashboards", 13, YELLOW, True)
mtxt(0.9, 3.95, 3.5, 0.8, [
    ("• NPS trend by region/model", LIGHT, False),
    ("• Customer health score heatmap", LIGHT, False),
    ("• Churn risk segmentation", LIGHT, False),
    ("• Battery degradation tracking", LIGHT, False),
], 10)

# Right path: Bedrock Agent
rect(5.0, 3.6, 4.3, 1.2, CARD, PURPLE, 1.5)
txt(5.3, 3.65, 3.8, 0.25, "Bedrock Agent (Claude)", 13, PURPLE, True)
mtxt(5.3, 3.95, 3.8, 0.8, [
    ("• Aurora pgvector for embeddings", LIGHT, False),
    ("• Natural language queries", LIGHT, False),
    ("• Cross-dataset synthesis", LIGHT, False),
    ("• Root cause analysis", LIGHT, False),
], 10)

# === RIGHT SIDE: Demo Scenario ===
txt(10.0, 1.4, 5.5, 0.3, "DEMO SCENARIO", 13, ORANGE, True)

# Question card
rect(10.0, 1.8, 5.5, 0.8, CARD, ORANGE, 2)
txt(10.3, 1.85, 5, 0.25, "User asks:", 11, MUTED)
txt(10.3, 2.1, 5, 0.4,
    "\"What's causing declining sentiment in the Southeast region?\"", 13, ORANGE, True)

arrow(12.75, 2.6, 12.75, 2.9, ORANGE, 1.5)

# Bedrock processing
rect(10.0, 2.9, 5.5, 1.4, CARD, PURPLE, 1.5)
txt(10.3, 2.95, 5, 0.25, "Bedrock Agent Reasoning", 12, PURPLE, True)
mtxt(10.3, 3.25, 5, 1.0, [
    ("1. Queries NPS data → declining in SE", LIGHT, False),
    ("2. Correlates with service records → battery", LIGHT, False),
    ("3. Checks telemetry → high temp exposure", LIGHT, False),
    ("4. Reviews contact center → wait times ↑", LIGHT, False),
], 10)

arrow(12.75, 4.3, 12.75, 4.55, PURPLE, 1.5)

# Answer card
rect(10.0, 4.55, 5.5, 1.6, CARD, GREEN, 2)
txt(10.3, 4.6, 5, 0.25, "AI-Synthesized Answer", 12, GREEN, True)
mtxt(10.3, 4.9, 5, 1.2, [
    ("Southeast NPS declined 12pts (Q3→Q4).", GREEN, True),
    ("Root causes identified:", WHITE, False),
    ("• Battery degradation in hot climates (42%)", LIGHT, False),
    ("• Service wait times up 3.2 days (31%)", LIGHT, False),
    ("• Delayed recall communication (27%)", LIGHT, False),
], 10)

# === BOTTOM: Key Metrics Bar ===
rect(0.6, 5.2, 8.7, 0.55, CARD, SKY, 1)
metrics = [
    ("Sub-second", "query response", 0.9),
    ("11 datasets", "integrated", 3.0),
    ("500K", "customers scored", 5.0),
    ("NL access", "for non-technical users", 6.8),
]
for val, label, x in metrics:
    txt(x, 5.25, 1.8, 0.2, val, 13, SKY, True)
    txt(x, 5.47, 1.8, 0.2, label, 9, MUTED)

# === BOTTOM RIGHT: Architecture callout ===
rect(10.0, 6.4, 5.5, 1.0, CARD, ORANGE, 1)
txt(10.3, 6.45, 5, 0.25, "Why This Matters", 12, ORANGE, True)
mtxt(10.3, 6.75, 5, 0.6, [
    ("Before: 5 teams, 5 dashboards, no shared truth", PINK, False),
    ("After: One question → AI synthesizes all sources", GREEN, False),
    ("Impact: Proactive retention, not reactive firefighting", ORANGE, False),
], 10)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide19_c360_deep_dive.pptx")
prs.save(out)
print(f"Saved: {out}")
