from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PINK = RGBColor(244, 114, 182)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(248, 113, 113)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def circ(s, l, t, sz, fill=CARD, border=None, bw=2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def pill(s, l, t, text, color, sz=12):
    pw = len(text) * 0.09 + 0.35
    ph = 0.32
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(pw), Inches(ph))
    dark = RGBColor(color[0] // 6, color[1] // 6, color[2] // 6)
    sh.fill.solid(); sh.fill.fore_color.rgb = dark
    sh.line.color.rgb = color; sh.line.width = Pt(1)
    tb = sh.text_frame; tb.word_wrap = False
    p = tb.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz); p.font.color.rgb = color
    return pw

# === TITLE ===
txt(slide, 0.6, 0.2, 14, 0.6, "CONNECTED VEHICLE PARTNER ECOSYSTEM", 36, ORANGE, True)
txt(slide, 0.6, 0.75, 14, 0.4,
    "Validated partner solutions at every layer — choose what works for your business", 20, LIGHT)

# ============================================================
# TOP SECTION — End-to-End CV Platforms
# ============================================================
txt(slide, 0.6, 1.4, 10, 0.4, "End-to-End Connected Vehicle Platforms", 20, WHITE, True)
txt(slide, 0.6, 1.8, 12, 0.3,
    "Turnkey platforms spanning the full stack — for OEMs who want a single partner", 13, MUTED)

platforms = [
    ("WirelessCar", GREEN,
     "11M+ vehicles in 100+ countries",
     "Full connectivity, journey intelligence, safety, security, and EV services"),
    ("Sibros", SKY,
     "Deep Connected Platform on AWS",
     "OTA updates, telematics, remote commands, and fleet-wide data management"),
    ("Excelfore", PURPLE,
     "20M+ vehicles, AWS Automotive Competency",
     "OTA, diagnostics, data aggregation, edge AI, and connected mobility"),
]

pw = 4.8
pgap = 0.2
px = 0.6
py = 2.2

for name, color, sub, desc in platforms:
    rect(slide, px, py, pw, 1.9, CARD, color, 1.5)

    # Partner name
    txt(slide, px + 0.25, py + 0.15, pw - 0.5, 0.35, name, 20, WHITE, True)

    # Subtitle
    txt(slide, px + 0.25, py + 0.55, pw - 0.5, 0.3, sub, 12, color)

    # Description
    txt(slide, px + 0.25, py + 0.95, pw - 0.5, 0.8, desc, 12, LIGHT)

    # "Full Stack" badge
    pill(slide, px + pw - 1.4, py + 1.5, "Full Stack", color, 10)

    px += pw + pgap

# Divider line
div_y = 4.35
line = slide.shapes.add_connector(1, Inches(0.6), Inches(div_y), Inches(15.4), Inches(div_y))
line.line.color.rgb = RGBColor(35, 45, 72); line.line.width = Pt(1)
line.line.dash_style = 2

# ============================================================
# BOTTOM SECTION — Platform Components
# ============================================================
txt(slide, 0.6, 4.5, 10, 0.4, "Platform Components", 20, WHITE, True)
txt(slide, 0.6, 4.9, 12, 0.3,
    "Best-of-breed partners at each layer — for OEMs building their own stack", 13, MUTED)

components = [
    ("Cellular Connectivity", SKY, ["Cubic Telecom", "KORE Wireless", "Eseye"]),
    ("Device Connectivity", GREEN, ["EMQX", "HiveMQ", "ClearBlade"]),
    ("Event Streaming", YELLOW, ["Confluent", "Redpanda"]),
    ("Stream Processing", ORANGE, ["Databricks", "Decodable", "Striim"]),
    ("Data & Analytics", PURPLE, ["Databricks", "Snowflake", "Teradata"]),
    ("Customer Engagement", PINK, ["Salesforce", "Genesys", "Zendesk"]),
    ("Observability", RED, ["Datadog", "Splunk", "Dynatrace"]),
]

sy = 5.3
rh = 0.42
rgap = 0.06

for i, (layer, color, partners) in enumerate(components):
    y = sy + i * (rh + rgap)

    # Row background
    rect(slide, 0.6, y, 14.8, rh, CARD)

    # Color dot + layer name
    circ(slide, 0.75, y + 0.06, 0.22, color, None)
    txt(slide, 1.05, y + 0.03, 3.5, 0.35, layer, 13, WHITE, True)

    # Partner pills
    ppx = 4.8
    for partner in partners:
        ppw = pill(slide, ppx, y + 0.05, partner, PURPLE, 11)
        ppx += ppw + 0.15

    # Marketplace badge
    pill(slide, 13.5, y + 0.05, "Marketplace", SKY, 9)

# === FOOTER ===
txt(slide, 0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide8_partners.pptx")
prs.save(out)
print(f"Saved: {out}")
