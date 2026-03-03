from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 33)

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
YELLOW = RGBColor(251, 191, 36)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)

def txt(l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def mtxt(l, t, w, h, lines, sz=10, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rect(l, t, w, h, fill=CARD, border=None, bw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def arrow(x1, y1, x2, y2, color=MUTED, width=1.5):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)

# === TITLE ===
txt(0.6, 0.2, 14, 0.5, "CUSTOMER 360 DEEP DIVE", 34, ORANGE, True)
txt(0.6, 0.65, 14, 0.35,
    "Unified customer profiles with AI-powered root cause analysis — 6 CDK stacks, 500K customers, 11 datasets", 17, LIGHT)

# =====================================================================
# LEFT: Architecture Pipeline (vertical flow)
# =====================================================================

# --- ① Data Lake ---
rect(0.6, 1.2, 7.3, 0.65, CARD, SKY, 1.5)
txt(0.9, 1.25, 3, 0.2, "① S3 Data Lake", 13, SKY, True)
txt(4.0, 1.25, 3.5, 0.2, "cx360-data-lake stack", 10, MUTED)
mtxt(0.9, 1.5, 6.5, 0.3, [
    ("Encrypted S3 bucket · Versioned · Access logging · Block public access · Partitioned: raw/ → processed/", LIGHT, False),
], 9)

arrow(4.3, 1.85, 4.3, 2.05, SKY, 1)

# --- ② Synthetic Data + Glue Catalog ---
rect(0.6, 2.05, 3.5, 1.0, CARD, GREEN, 1.5)
txt(0.9, 2.1, 3, 0.2, "② Synthetic Data", 12, GREEN, True)
mtxt(0.9, 2.35, 3, 0.65, [
    ("500K customers · 200 dealers", GREEN, False),
    ("Faker + Aurora PostgreSQL", LIGHT, False),
    ("10 years of history (2015-2024)", MUTED, False),
], 9)

rect(4.4, 2.05, 3.5, 1.0, CARD, GREEN, 1.5)
txt(4.7, 2.1, 3, 0.2, "② Glue Catalog", 12, GREEN, True)
mtxt(4.7, 2.35, 3, 0.65, [
    ("11 tables: accounts, vehicles,", LIGHT, False),
    ("cases, surveys, service_appts,", LIGHT, False),
    ("opportunities, dealers, users...", MUTED, False),
], 9)

arrow(4.3, 3.05, 4.3, 3.25, GREEN, 1)

# --- ③ ETL + Health Scores ---
rect(0.6, 3.25, 7.3, 1.15, CARD, PURPLE, 1.5)
txt(0.9, 3.3, 3, 0.2, "③ Glue ETL Jobs", 13, PURPLE, True)
txt(4.0, 3.3, 3.5, 0.2, "cx360-etl stack", 10, MUTED)
mtxt(0.9, 3.55, 3.3, 0.8, [
    ("process_customer_360.py:", WHITE, True),
    ("Join accounts + vehicles + cases", LIGHT, False),
    ("+ surveys + service appointments", LIGHT, False),
    ("→ processed/customer-360/ (Parquet)", MUTED, False),
], 9)
mtxt(4.3, 3.55, 3.3, 0.8, [
    ("calculate_health_scores.py:", WHITE, True),
    ("Weighted score (0-100):", LIGHT, False),
    ("Recency 30% · Satisfaction 25%", LIGHT, False),
    ("Support 20% · Engage 15% · Pay 10%", MUTED, False),
], 9)

# Split arrows to Athena/QuickSight and Bedrock
arrow(2.5, 4.4, 2.5, 4.65, PURPLE, 1)
arrow(5.8, 4.4, 5.8, 4.65, PURPLE, 1)

# --- ④ Athena + QuickSight (left) ---
rect(0.6, 4.65, 3.5, 1.3, CARD, YELLOW, 1.5)
txt(0.9, 4.7, 3, 0.2, "④ Athena + QuickSight", 12, YELLOW, True)
mtxt(0.9, 4.95, 3, 1.0, [
    ("Athena workgroup + 7 views:", WHITE, False),
    ("• customer_health_scores", LIGHT, False),
    ("• at_risk_revenue", LIGHT, False),
    ("• kpi_trends, operational_kpis", LIGHT, False),
    ("• issue_categories, revenue_breakdown", MUTED, False),
    ("QuickSight: OEM Business Overview", YELLOW, True),
], 9)

# --- ④ Bedrock Agent (right) ---
rect(4.4, 4.65, 3.5, 1.3, CARD, PURPLE, 1.5)
txt(4.7, 4.7, 3, 0.2, "④ Bedrock Agent", 12, PURPLE, True)
mtxt(4.7, 4.95, 3, 1.0, [
    ("Knowledge Base:", WHITE, True),
    ("• S3 → Titan Embed v1 (1536 dims)", LIGHT, False),
    ("• Aurora pgvector storage", LIGHT, False),
    ("• Playbooks: battery, churn, analytics", MUTED, False),
    ("Action Group (Lambda):", WHITE, True),
    ("6 API paths for Athena queries", PURPLE, False),
], 9)

# =====================================================================
# RIGHT: Demo + Agent Details
# =====================================================================

txt(8.3, 1.2, 7, 0.25, "BEDROCK AGENT ACTION GROUP", 13, ORANGE, True)

rect(8.3, 1.5, 7.1, 2.1, CARD, PURPLE, 1.5)
txt(8.6, 1.55, 6.5, 0.2, "Lambda handles 6 API paths — each queries Athena views:", 10, LIGHT)

mtxt(8.6, 1.8, 3.2, 1.7, [
    ("/query-health-segments", PURPLE, True),
    ("  Customer health by segment", LIGHT, False),
    ("/query-at-risk-customers", PURPLE, True),
    ("  Revenue at risk + churn prob", LIGHT, False),
    ("/query-sentiment-trends", PURPLE, True),
    ("  NPS trends over time", LIGHT, False),
], 9)

mtxt(12.0, 1.8, 3.2, 1.7, [
    ("/query-root-causes", PURPLE, True),
    ("  Issue categories + root causes", LIGHT, False),
    ("/query-customer-360", PURPLE, True),
    ("  Full customer profile lookup", LIGHT, False),
    ("/query-dashboard-metrics", PURPLE, True),
    ("  Operational KPIs + trends", LIGHT, False),
], 9)

# --- Demo Scenario ---
txt(8.3, 3.8, 7, 0.25, "DEMO SCENARIO", 13, ORANGE, True)

rect(8.3, 4.1, 7.1, 0.65, CARD, ORANGE, 2)
txt(8.6, 4.15, 6.5, 0.2, "User asks:", 10, MUTED)
txt(8.6, 4.35, 6.5, 0.3,
    "\"What's causing declining sentiment in the Southeast region?\"", 13, ORANGE, True)

arrow(11.85, 4.75, 11.85, 5.0, ORANGE, 1.5)

rect(8.3, 5.0, 7.1, 1.3, CARD, PURPLE, 1.5)
txt(8.6, 5.05, 6.5, 0.2, "Agent Reasoning Chain", 12, PURPLE, True)
mtxt(8.6, 5.3, 6.5, 0.95, [
    ("1. /query-sentiment-trends → NPS declining 12pts in SE (Q3→Q4)", LIGHT, False),
    ("2. /query-root-causes → battery issues 42%, service delays 31%, recall comms 27%", LIGHT, False),
    ("3. Knowledge Base retrieval → battery-remediation-playbook.md", LIGHT, False),
    ("4. /query-at-risk-customers → 2,340 high-value customers at risk ($4.2M revenue)", LIGHT, False),
], 9)

arrow(11.85, 6.3, 11.85, 6.55, PURPLE, 1.5)

rect(8.3, 6.55, 7.1, 1.0, CARD, GREEN, 2)
txt(8.6, 6.6, 6.5, 0.2, "Synthesized Answer", 12, GREEN, True)
mtxt(8.6, 6.85, 6.5, 0.65, [
    ("\"SE NPS dropped 12pts. Root causes: battery degradation in hot climates (42%),", GREEN, False),
    ("service wait times up 3.2 days (31%), delayed recall comms (27%).", LIGHT, False),
    ("Recommended: Proactive battery replacement per playbook. 2,340 customers at risk.\"", LIGHT, False),
], 9)

# === BOTTOM: Key Metrics Bar ===
rect(0.5, 7.85, 15.0, 0.45, CARD, ORANGE, 1)
metrics = [
    ("6 CDK stacks", "modular deployment", 0.8),
    ("500K customers", "synthetic dataset", 3.0),
    ("11 Glue tables", "unified catalog", 5.2),
    ("7 Athena views", "analytical layer", 7.4),
    ("Titan Embed v1", "1536-dim vectors", 9.6),
    ("6 API paths", "agent action group", 11.8),
    ("pgvector", "Aurora PostgreSQL", 13.8),
]
for val, label, x in metrics:
    txt(x, 7.88, 2, 0.2, val, 10, ORANGE, True)
    txt(x, 8.07, 2, 0.2, label, 8, MUTED)

# === FOOTER ===
txt(0.6, 8.6, 12, 0.3,
    "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved.", 10, MUTED)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slide19_c360_deep_dive_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
