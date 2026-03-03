from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = '/Users/givenand/Downloads/2026_aws_powerpoint_template_v1_07162eba (1).pptx'
OUTDIR = os.path.dirname(os.path.abspath(__file__))

ORANGE = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(190, 200, 220)
MUTED = RGBColor(110, 125, 150)
SKY = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(244, 114, 182)
CARD = RGBColor(20, 28, 52)
DARK = RGBColor(14, 18, 32)
DKBLUE = RGBColor(35, 47, 62)  # AWS dark blue

# ============================================================
# VERSION A — Two Content layout (layout 24)
# ============================================================
prs_a = Presentation(TEMPLATE)
# Delete all existing slides
while len(prs_a.slides) > 0:
    rId = prs_a.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs_a.part.drop_rel(rId)
    prs_a.slides._sldIdLst.remove(prs_a.slides._sldIdLst[0])

layout_a = prs_a.slide_layouts[24]  # Two Content
slide_a = prs_a.slides.add_slide(layout_a)

# Title
title = slide_a.placeholders[0]
title.text = "Automotive Data Platform — Architecture Overview"

# Left placeholder — four pillars
left = slide_a.placeholders[1]
tf = left.text_frame
tf.clear()

pillars = [
    ("Customer 360 Analytics", "Unified customer profiles with AI-powered insights via Amazon Bedrock agents and QuickSight dashboards"),
    ("Predictive Maintenance", "ML-based tire failure prediction 7-14 days in advance using SageMaker Random Cut Forest"),
    ("Data Mesh Foundation", "Domain-oriented data ownership with SageMaker Unified Studio and DataZone catalog"),
    ("Data Governance", "Fine-grained access control with Lake Formation, automated PII detection with Macie"),
]

for i, (name, desc) in enumerate(pillars):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_before = Pt(12) if i > 0 else Pt(0)

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(51, 51, 51)
    p2.space_after = Pt(4)

# Right placeholder — six layers
right = slide_a.placeholders[2]
tf2 = right.text_frame
tf2.clear()

layers = [
    ("1. Data Sources", "CRM, DMS, telemetry, service, contact center, third-party"),
    ("2. Data Ingestion", "IoT Core (MQTT), MSK (Kafka), Kinesis, API Gateway"),
    ("3. Storage & Governance", "S3 data lake, Glue Catalog, Lake Formation, DataZone"),
    ("4. Processing", "Managed Flink (real-time), Glue ETL (batch), SageMaker (ML)"),
    ("5. Analytics", "Athena (SQL), QuickSight (dashboards), Bedrock (AI agents)"),
    ("6. Action", "Amazon Connect (proactive outreach), EventBridge (alerts)"),
]

for i, (layer, services) in enumerate(layers):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = layer
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(35, 47, 62)
    p.space_before = Pt(8) if i > 0 else Pt(0)

    p2 = tf2.add_paragraph()
    p2.text = services
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.space_after = Pt(2)

prs_a.save(os.path.join(OUTDIR, "slide18_adp_arch_v_a.pptx"))
print("Saved: slide18_adp_arch_v_a.pptx (Two Content layout)")

# ============================================================
# VERSION B — Title Only layout (layout 12) with custom shapes
# ============================================================
prs_b = Presentation(TEMPLATE)
while len(prs_b.slides) > 0:
    rId = prs_b.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs_b.part.drop_rel(rId)
    prs_b.slides._sldIdLst.remove(prs_b.slides._sldIdLst[0])

layout_b = prs_b.slide_layouts[12]  # Title Only
slide_b = prs_b.slides.add_slide(layout_b)

# Title
slide_b.placeholders[0].text = "Automotive Data Platform — Architecture Overview"

def txt(s, l, t, w, h, text, sz=14, color=RGBColor(51,51,51), bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold

def rect(s, l, t, w, h, fill=None, border=None, bw=1.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()

# Subtitle
txt(slide_b, 0.5, 0.95, 12, 0.3,
    "Four integrated solutions on a six-layer serverless architecture", 13, RGBColor(100,100,100))

# === Four pillar cards across the top ===
pillars_b = [
    ("Customer 360\nAnalytics", "Bedrock Agent\nQuickSight", SKY),
    ("Predictive\nMaintenance", "SageMaker\nIoT Core", GREEN),
    ("Data Mesh\nFoundation", "Unified Studio\nDataZone", PURPLE),
    ("Data\nGovernance", "Lake Formation\nMacie", PINK),
]

pw = 2.9
pgap = 0.2
px = 0.5
py = 1.45

for name, services, color in pillars_b:
    rect(slide_b, px, py, pw, 1.5, RGBColor(245,247,250), color, 2)
    txt(slide_b, px + 0.15, py + 0.1, pw - 0.3, 0.6, name, 13, DKBLUE, True, PP_ALIGN.CENTER)
    txt(slide_b, px + 0.15, py + 0.8, pw - 0.3, 0.55, services, 10, RGBColor(100,100,100), False, PP_ALIGN.CENTER)
    px += pw + pgap

# === Six-layer architecture stack ===
layers_b = [
    ("Data Sources", "CRM · DMS · Telemetry · Service · Contact Center · Third-Party", ORANGE),
    ("Data Ingestion", "IoT Core (MQTT) · Amazon MSK (Kafka) · Kinesis · API Gateway", SKY),
    ("Storage & Governance", "S3 Data Lake · Glue Catalog · Lake Formation · DataZone", GREEN),
    ("Processing", "Managed Flink (real-time) · Glue ETL (batch) · SageMaker (ML)", PURPLE),
    ("Analytics & AI", "Athena (SQL) · QuickSight (dashboards) · Bedrock (AI agents)", RGBColor(251, 191, 36)),
    ("Action", "Amazon Connect (proactive outreach) · EventBridge (alerts) · SNS", PINK),
]

ly = 3.2
lh = 0.55
lgap = 0.08

for name, services, color in layers_b:
    # Color bar on left
    sh = slide_b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(ly), Inches(0.08), Inches(lh))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

    # Row background
    rect(slide_b, 0.65, ly, 12.15, lh, RGBColor(245,247,250))

    # Layer name
    txt(slide_b, 0.75, ly + 0.05, 2.2, 0.4, name, 12, DKBLUE, True)

    # Services
    txt(slide_b, 3.0, ly + 0.05, 9.5, 0.4, services, 11, RGBColor(80,80,80))

    # Down arrow between layers (except last)
    if name != "Action":
        arr_x = 6.5
        arr = slide_b.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(arr_x), Inches(ly + lh), Inches(0.2), Inches(lgap))
        arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(200, 200, 200)
        arr.line.fill.background()

    ly += lh + lgap

# Key stats bar
sy = ly + 0.2
rect(slide_b, 0.5, sy, 12.3, 0.5, RGBColor(235,240,245))
stats = [
    ("500K customers", 0.7),
    ("1.4M interactions", 3.2),
    ("900K service records", 5.7),
    ("11 datasets", 8.2),
    ("~$114/mo", 10.0),
    ("MIT-0 License", 11.3),
]
for label, x in stats:
    txt(slide_b, x, sy + 0.08, 2, 0.35, label, 11, DKBLUE, True, PP_ALIGN.LEFT)

prs_b.save(os.path.join(OUTDIR, "slide18_adp_arch_v_b.pptx"))
print("Saved: slide18_adp_arch_v_b.pptx (Title Only + custom shapes)")
